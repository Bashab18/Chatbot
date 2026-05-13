from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x11, 0x17)
ACCENT   = RGBColor(0x10, 0xB9, 0x81)   # green
BLUE     = RGBColor(0x3B, 0x82, 0xF6)
PURPLE   = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
PINK     = RGBColor(0xEC, 0x48, 0x99)
TEAL     = RGBColor(0x14, 0xB8, 0xA6)
WHITE    = RGBColor(0xF1, 0xF5, 0xF9)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
DIM      = RGBColor(0x47, 0x55, 0x69)
CARD_BG  = RGBColor(0x1E, 0x22, 0x2E)
BORDER   = RGBColor(0x1E, 0x29, 0x3B)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txbox(slide, text, left, top, width, height,
          size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return tb

def rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=BORDER, line_width=Pt(0.75)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = line_width
    else:
        shape.line.fill.background()
    return shape

def label_tag(slide, text, left, top):
    txbox(slide, text, left, top, Inches(6), Inches(0.3),
          size=9, bold=True, color=ACCENT)

def heading(slide, text, left, top, width=Inches(12), size=32):
    txbox(slide, text, left, top, width, Inches(0.6),
          size=size, bold=True, color=WHITE)

def subtext(slide, text, left, top, width=Inches(9), size=13):
    txbox(slide, text, left, top, width, Inches(0.8),
          size=size, color=MUTED)

def bullet_block(slide, title, items, left, top, width, height, title_color=ACCENT):
    r = rect(slide, left, top, width, height)
    txbox(slide, title, left + Inches(0.18), top + Inches(0.15),
          width - Inches(0.2), Inches(0.3), size=10, bold=True, color=title_color)
    y = top + Inches(0.48)
    row_h = (height - Inches(0.55)) / max(len(items), 1)
    for item in items:
        txbox(slide, f"→  {item}", left + Inches(0.18), y,
              width - Inches(0.25), row_h + Inches(0.05),
              size=11, color=MUTED)
        y += row_h

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
bg(s1)

# Green accent bar left edge
rect(s1, Inches(0), Inches(2.8), Inches(0.06), Inches(2.2),
     fill_color=ACCENT, line_color=None)

# Logo box
logo = rect(s1, Inches(1.2), Inches(1.2), Inches(0.72), Inches(0.72),
            fill_color=ACCENT, line_color=None)
txbox(s1, "✦", Inches(1.2), Inches(1.2), Inches(0.72), Inches(0.72),
      size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

label_tag(s1, "SYSTEM DESIGN & ARCHITECTURE", Inches(1.2), Inches(2.1))

# Main title
txbox(s1, "CIRA Health &", Inches(1.2), Inches(2.55), Inches(8), Inches(0.9),
      size=48, bold=True, color=WHITE)
txbox(s1, "Wellness Chatbot", Inches(1.2), Inches(3.35), Inches(8), Inches(0.9),
      size=48, bold=True, color=ACCENT)

txbox(s1, "A full-stack AI-powered health assistant built on Gemini LLMs,\n"
          "BM25 Retrieval-Augmented Generation, and a SQLite persistence layer\n"
          "— deployed across Netlify and Render.",
      Inches(1.2), Inches(4.35), Inches(9), Inches(1.0), size=14, color=MUTED)

# Pill row
pills = [
    ("React 18", ACCENT), ("Node.js / Express", ACCENT), ("Google Gemini API", ACCENT),
    ("SQLite + BM25", DIM), ("ElevenLabs TTS", DIM), ("JWT Auth", DIM),
    ("Netlify + Render", DIM),
]
px = Inches(1.2)
for txt, col in pills:
    w = Inches(len(txt) * 0.095 + 0.35)
    r = rect(s1, px, Inches(5.55), w, Inches(0.33),
             fill_color=RGBColor(0x1A, 0x1E, 0x28), line_color=col)
    txbox(s1, txt, px + Inches(0.12), Inches(5.55), w - Inches(0.1), Inches(0.33),
          size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    px += w + Inches(0.12)

# Slide number
txbox(s1, "1 / 5", Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
      size=10, color=DIM, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — High-Level Architecture
# ═══════════════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
bg(s2)

label_tag(s2, "SLIDE 2 OF 5", Inches(0.5), Inches(0.3))
heading(s2, "High-Level Architecture", Inches(0.5), Inches(0.65))
subtext(s2, "Three-tier architecture: React SPA frontend · Express REST API backend · Hybrid SQLite + in-memory RAG store",
        Inches(0.5), Inches(1.2), width=Inches(12.3))

cols = [
    ("FRONTEND — Netlify CDN", ACCENT, [
        ("React 18 SPA", "React Router v7, Context API for auth state"),
        ("Chat Interface",  "Streaming UI, Markdown rendering, suggestion chips"),
        ("Admin Panel",     "Dashboard, KB management, user controls, settings"),
        ("TTS Playback",    "ElevenLabs audio via /api/tts, Web Audio API"),
        ("Profile & Memory","User health profile, AI-extracted memory cards"),
    ]),
    ("BACKEND — Render (Node 20)", BLUE, [
        ("Express REST API",   "20+ endpoints, JWT middleware, role-based guards"),
        ("Chat Engine",        "Gemini multi-turn chat, system prompt, style control"),
        ("RAG Pipeline",       "BM25 tokenisation, chunk search, context assembly"),
        ("Document Ingestion", "PDF parsing, text chunking, BM25 index building"),
        ("Memory Extraction",  "Gemini analyses history, saves user facts to DB"),
    ]),
    ("EXTERNAL SERVICES", PURPLE, [
        ("Google Gemini API", "Gemini 3.x / 2.x / Gemma 4 — chat & title gen"),
        ("ElevenLabs TTS",   "Voice synthesis, Turbo v2.5, per-message audio"),
        ("Google Fit OAuth",  "Fitness snapshot (steps, HR, weight) via OAuth2"),
        ("Google Search",     "Live web grounding — optional, weight-controlled"),
        ("Netlify CDN",       "Global static hosting with /api/* proxy to Render"),
    ]),
]

col_w = Inches(4.1)
for i, (title, color, items) in enumerate(cols):
    cx = Inches(0.4) + i * (col_w + Inches(0.18))
    cy = Inches(1.85)
    ch = Inches(5.3)
    rect(s2, cx, cy, col_w, ch)

    # Column header
    txbox(s2, title, cx + Inches(0.18), cy + Inches(0.15),
          col_w - Inches(0.2), Inches(0.3), size=9, bold=True, color=color)

    # Items
    iy = cy + Inches(0.55)
    item_h = Inches(0.85)
    for (name, desc) in items:
        txbox(s2, f"● {name}", cx + Inches(0.18), iy,
              col_w - Inches(0.25), Inches(0.25), size=11, bold=True, color=WHITE)
        txbox(s2, desc, cx + Inches(0.28), iy + Inches(0.26),
              col_w - Inches(0.35), Inches(0.52), size=10, color=MUTED)
        iy += item_h

txbox(s2, "2 / 5", Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
      size=10, color=DIM, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — RAG Pipeline
# ═══════════════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
bg(s3)

label_tag(s3, "SLIDE 3 OF 5", Inches(0.5), Inches(0.3))
heading(s3, "RAG Pipeline — Knowledge Retrieval", Inches(0.5), Inches(0.65))
subtext(s3, "BM25-based retrieval augments every Gemini response with relevant excerpts from the admin-curated knowledge base.",
        Inches(0.5), Inches(1.2), width=Inches(12.3))

steps = [
    ("📄", "1. Upload",    "Admin uploads PDF/TXT. pdf-parse extracts raw text."),
    ("✂️", "2. Chunk",     "Split into overlapping ~800-token chunks preserving sentence boundaries."),
    ("🔤", "3. Tokenise",  "Each chunk tokenised into word arrays, stored as JSON in SQLite."),
    ("🔍", "4. Rank",      "User query tokenised. BM25 scores all chunks; top-K returned."),
    ("🤖", "5. Inject",    "Top chunks prepended to system prompt. Gemini responds grounded."),
]

step_w = Inches(2.4)
step_h = Inches(2.1)
for i, (icon, title, desc) in enumerate(steps):
    sx = Inches(0.35) + i * (step_w + Inches(0.1))
    sy = Inches(1.85)
    rect(s3, sx, sy, step_w, step_h)
    txbox(s3, icon, sx, sy + Inches(0.15), step_w, Inches(0.45),
          size=22, align=PP_ALIGN.CENTER)
    txbox(s3, title, sx, sy + Inches(0.62), step_w, Inches(0.3),
          size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txbox(s3, desc, sx + Inches(0.12), sy + Inches(0.95),
          step_w - Inches(0.2), Inches(1.0), size=10, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        txbox(s3, "›", sx + step_w + Inches(0.01), sy + Inches(0.8),
              Inches(0.12), Inches(0.4), size=18, color=DIM, align=PP_ALIGN.CENTER)

# Two detail boxes
bx, by, bw, bh = Inches(0.35), Inches(4.15), Inches(6.1), Inches(2.9)
bullet_block(s3, "Retrieval Controls (Admin-tunable)", [
    "ragTopK — number of chunks to retrieve (default 5)",
    "ragMinScore — BM25 minimum relevance threshold",
    "ragWeight — blend KB vs. AI own knowledge (0–100)",
    "webSearchWeight — live Google Search grounding",
    "Strict KB-only mode refuses off-topic questions",
], bx, by, bw, bh)

bullet_block(s3, "Profile Context Bypass", [
    "Users with health profiles always get a response",
    "Profile data injected even when KB has no hits",
    "AI memories extracted from past conversations",
    "Google Fit snapshot (steps, HR, weight) injected",
    "Personalised answers without KB dependency",
], bx + bw + Inches(0.18), by, bw, bh, title_color=BLUE)

txbox(s3, "3 / 5", Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
      size=10, color=DIM, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Data Layer
# ═══════════════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
bg(s4)

label_tag(s4, "SLIDE 4 OF 5", Inches(0.5), Inches(0.3))
heading(s4, "Data Layer — SQLite Schema", Inches(0.5), Inches(0.65))
subtext(s4, "A single SQLite database (WAL mode) persists all state — users, conversations, KB documents, settings, and AI memories.",
        Inches(0.5), Inches(1.2), width=Inches(12.3))

tables = [
    ("👤 users",         ACCENT,  ["id (TEXT PK)", "email (UNIQUE)", "name, role", "password_hash (bcrypt)", "login_count, last_login"]),
    ("💬 conversations", BLUE,    ["id (TEXT PK)", "user_id (FK → users)", "title (auto-generated)", "created_at, updated_at"]),
    ("📝 messages",      PURPLE,  ["id (TEXT PK)", "conversation_id (FK)", "role: user | assistant", "text (TEXT)", "refs (JSON sources)"]),
    ("📚 kb_documents",  ORANGE,  ["id (TEXT PK)", "name (filename)", "chunk_count (INTEGER)", "added_at (epoch ms)"]),
    ("🔤 kb_chunks",     TEAL,    ["doc_id (FK → kb_docs)", "doc_name (TEXT)", "text (raw chunk)", "embedding (JSON tokens[])"]),
    ("⚙️ settings+",    PINK,    ["settings (key, value)", "user_profiles (health)", "user_memories (AI facts)", "fitness_snapshots (Fit)"]),
]

tw = Inches(4.1)
th = Inches(2.5)
for i, (name, color, fields) in enumerate(tables):
    col = i % 3
    row = i // 3
    tx = Inches(0.35) + col * (tw + Inches(0.18))
    ty = Inches(1.85) + row * (th + Inches(0.14))

    rect(s4, tx, ty, tw, th)
    # Header bar
    hdr = rect(s4, tx, ty, tw, Inches(0.38), fill_color=color, line_color=None)
    txbox(s4, name, tx + Inches(0.12), ty + Inches(0.05),
          tw - Inches(0.15), Inches(0.28), size=11, bold=True, color=WHITE)

    fy = ty + Inches(0.44)
    fh = (th - Inches(0.48)) / max(len(fields), 1)
    for field in fields:
        txbox(s4, field, tx + Inches(0.15), fy,
              tw - Inches(0.2), fh, size=10, color=MUTED)
        fy += fh

# Note bar
note_y = Inches(6.45)
rect(s4, Inches(0.35), note_y, Inches(12.6), Inches(0.72),
     fill_color=RGBColor(0x0D, 0x1F, 0x18), line_color=RGBColor(0x10, 0x4D, 0x38))
txbox(s4, "Persistence on Render:  SQLite stored on a 1 GB persistent disk at "
          "/opt/render/project/src/backend/data — survives restarts & redeploys. "
          "KB chunks loaded into memory on startup for fast BM25 scoring.",
      Inches(0.55), note_y + Inches(0.08), Inches(12.2), Inches(0.6),
      size=11, color=MUTED)

txbox(s4, "4 / 5", Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
      size=10, color=DIM, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Deployment & API
# ═══════════════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
bg(s5)

label_tag(s5, "SLIDE 5 OF 5", Inches(0.5), Inches(0.3))
heading(s5, "Deployment & API Surface", Inches(0.5), Inches(0.65))
subtext(s5, "React SPA on Netlify CDN · Express API on Render with persistent disk · /api/* requests proxied transparently.",
        Inches(0.5), Inches(1.2), width=Inches(12.3))

# Left column – two stacked cards
lx, ly, lw = Inches(0.35), Inches(1.85), Inches(5.6)

# Frontend card
feh = Inches(2.4)
rect(s5, lx, ly, lw, feh)
txbox(s5, "🌐  Frontend — Netlify", lx + Inches(0.18), ly + Inches(0.12),
      lw - Inches(0.2), Inches(0.3), size=12, bold=True, color=ACCENT)
fe_rows = [
    ("Host",      "chatbotmhealth.netlify.app"),
    ("Build",     "npm run build → /build"),
    ("Routing",   "SPA fallback → index.html"),
    ("API proxy", "/api/* → Render backend"),
    ("Stack",     "React 18 + React Router v7"),
]
ry = ly + Inches(0.5)
for k, v in fe_rows:
    txbox(s5, k, lx + Inches(0.18), ry, Inches(1.5), Inches(0.28), size=10, color=DIM)
    txbox(s5, v, lx + Inches(1.8),  ry, lw - Inches(1.9), Inches(0.28), size=10, color=WHITE, bold=True)
    ry += Inches(0.33)

# Backend card
bey = ly + feh + Inches(0.14)
beh = Inches(2.55)
rect(s5, lx, bey, lw, beh)
txbox(s5, "⚙️  Backend — Render", lx + Inches(0.18), bey + Inches(0.12),
      lw - Inches(0.2), Inches(0.3), size=12, bold=True, color=BLUE)
be_rows = [
    ("Host",    "chatbot-bw9p.onrender.com"),
    ("Runtime", "Node 20.x"),
    ("Start",   "node server.js"),
    ("Disk",    "1 GB persistent (SQLite)"),
    ("Auth",    "JWT HS256 + bcrypt passwords"),
]
ry = bey + Inches(0.5)
for k, v in be_rows:
    txbox(s5, k, lx + Inches(0.18), ry, Inches(1.5), Inches(0.28), size=10, color=DIM)
    txbox(s5, v, lx + Inches(1.8),  ry, lw - Inches(1.9), Inches(0.28), size=10, color=WHITE, bold=True)
    ry += Inches(0.38)

# Right column – API endpoints
rx, ry_start = Inches(6.15), Inches(1.85)
rw = Inches(6.85)
rh = Inches(5.35)
rect(s5, rx, ry_start, rw, rh)
txbox(s5, "Key API Endpoints", rx + Inches(0.18), ry_start + Inches(0.15),
      rw - Inches(0.2), Inches(0.3), size=12, bold=True, color=WHITE)

endpoints = [
    ("POST", "/api/auth/signup  ·  /api/auth/login",   ACCENT),
    ("POST", "/api/chat",                               ACCENT),
    ("POST", "/api/tts",                                ACCENT),
    ("POST", "/api/title",                              ACCENT),
    ("GET",  "/api/conversations",                      BLUE),
    ("POST", "/api/kb/upload",                          ACCENT),
    ("DEL",  "/api/kb/:id",                             RGBColor(0xEF,0x44,0x44)),
    ("GET",  "/api/admin/stats  ·  /api/admin/users",  BLUE),
    ("GET",  "/api/user/profile  ·  /api/user/memories", BLUE),
    ("GET",  "/api/fitness/callback  (OAuth2)",         PURPLE),
]

method_colors = {"POST": ACCENT, "GET": BLUE, "DEL": RGBColor(0xEF,0x44,0x44)}
ey = ry_start + Inches(0.55)
for method, path, _ in endpoints:
    mc = method_colors.get(method, MUTED)
    # method pill
    mp = rect(s5, rx + Inches(0.18), ey, Inches(0.55), Inches(0.28),
              fill_color=RGBColor(0x1A,0x2A,0x22) if mc == ACCENT
              else RGBColor(0x1A,0x22,0x35) if mc == BLUE
              else RGBColor(0x2A,0x1A,0x1A),
              line_color=mc)
    txbox(s5, method, rx + Inches(0.18), ey, Inches(0.55), Inches(0.28),
          size=8, bold=True, color=mc, align=PP_ALIGN.CENTER)
    txbox(s5, path, rx + Inches(0.82), ey, rw - Inches(0.95), Inches(0.28),
          size=10, color=MUTED)
    ey += Inches(0.46)

txbox(s5, "5 / 5", Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
      size=10, color=DIM, align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────
prs.save("/home/user/Chatbot/CIRA_Architecture.pptx")
print("Saved: CIRA_Architecture.pptx")
