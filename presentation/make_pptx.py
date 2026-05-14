#!/usr/bin/env python3
"""CIRA Chatbot – 70-Slide Conference Presentation  (v2 – improved)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ══════════════════════════════════════════════════════════════════
#  COLOUR PALETTE
# ══════════════════════════════════════════════════════════════════
BG      = RGBColor(0x08, 0x0d, 0x18)
BG2     = RGBColor(0x0b, 0x12, 0x20)
BG3     = RGBColor(0x0f, 0x17, 0x2a)
CARD    = RGBColor(0x10, 0x1c, 0x32)
CARD2   = RGBColor(0x15, 0x22, 0x3c)
STROKE  = RGBColor(0x1e, 0x2e, 0x4a)
BLUE    = RGBColor(0x4f, 0x9c, 0xf9)
PURPLE  = RGBColor(0xa8, 0x55, 0xf7)
CYAN    = RGBColor(0x22, 0xd3, 0xee)
GREEN   = RGBColor(0x34, 0xd3, 0x99)
YELLOW  = RGBColor(0xfb, 0xbf, 0x24)
RED     = RGBColor(0xf8, 0x71, 0x71)
WHITE   = RGBColor(0xff, 0xff, 0xff)
LIGHT   = RGBColor(0xe2, 0xe8, 0xf0)
MUTED   = RGBColor(0x94, 0xa3, 0xb8)
DIM     = RGBColor(0x2e, 0x3e, 0x58)
CODBG   = RGBColor(0x0d, 0x11, 0x17)
CODBAR  = RGBColor(0x16, 0x1b, 0x27)
CODBDR  = RGBColor(0x21, 0x2d, 0x45)
CODTXT  = RGBColor(0xa5, 0xd6, 0xff)

def _dim(c, f=6):
    return RGBColor(min(255,c[0]//f+12), min(255,c[1]//f+12), min(255,c[2]//f+12))

# ══════════════════════════════════════════════════════════════════
#  PRIMITIVES
# ══════════════════════════════════════════════════════════════════
def new_slide(prs, bg=BG):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    f = sl.background.fill; f.solid(); f.fore_color.rgb = bg
    return sl

def T(sl, text, l, t, w, h, size=13, bold=False, color=LIGHT,
      align=PP_ALIGN.LEFT, italic=False):
    bx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    bx.word_wrap = True
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = 'Calibri'
    return bx

def R(sl, l, t, w, h, fill=CARD, line=None, lw=0.75, radius=False):
    sh = sl.shapes.add_shape(5 if radius else 1,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    return sh

def hline(sl, l, t, w, color=STROKE):
    R(sl, l, t, w, 0.012, fill=color)

def snum(sl, n, total=70):
    T(sl, f'{n:02d} / {total}', 12.0, 7.18, 1.25, 0.25,
      size=8, color=DIM, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
#  BADGE
# ══════════════════════════════════════════════════════════════════
def badge(sl, text, l, t, accent=BLUE):
    w = max(1.4, len(text) * 0.107 + 0.6)
    sh = R(sl, l, t, w, 0.28, fill=_dim(accent, 5), line=accent, lw=0.75, radius=True)
    tf = sh.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text.upper()
    r.font.size = Pt(7.5); r.font.bold = True
    r.font.color.rgb = accent; r.font.name = 'Calibri'

# ══════════════════════════════════════════════════════════════════
#  SLIDE HEADER  (badge + title + divider)
# ══════════════════════════════════════════════════════════════════
def slide_header(sl, badge_txt, title, accent=BLUE):
    badge(sl, badge_txt, 0.48, 0.32, accent)
    T(sl, title, 0.48, 0.67, 12.35, 0.68, size=27, bold=True, color=WHITE)
    hline(sl, 0.48, 1.4, 12.37)

# ══════════════════════════════════════════════════════════════════
#  CARD  (left accent strip + icon + title + description)
# ══════════════════════════════════════════════════════════════════
def card(sl, l, t, w, h, icon, title, desc, accent=BLUE):
    R(sl, l, t, w, h, fill=_dim(accent, 9), line=_dim(accent, 4), lw=0.6, radius=True)
    R(sl, l, t, 0.06, h, fill=accent)          # left accent strip
    T(sl, icon,  l+0.15, t+0.10, 0.5,   0.48, size=18)
    T(sl, title, l+0.15, t+0.58, w-0.27, 0.30, size=11.5, bold=True, color=WHITE)
    T(sl, desc,  l+0.15, t+0.92, w-0.27, max(0.3, h-1.05), size=9.5, color=MUTED)

# ══════════════════════════════════════════════════════════════════
#  ARCH BOX  (left accent + title + description)
# ══════════════════════════════════════════════════════════════════
def arch_box(sl, l, t, w, icon, title, desc, accent=BLUE, h=0.84):
    R(sl, l, t, w, h, fill=CARD, line=STROKE, lw=0.5)
    R(sl, l, t, 0.055, h, fill=accent)
    T(sl, f'{icon}  {title}', l+0.16, t+0.07, w-0.24, 0.30, size=11.5, bold=True, color=WHITE)
    T(sl, desc, l+0.16, t+0.40, w-0.24, max(0.22, h-0.50), size=9.5, color=MUTED)

# ══════════════════════════════════════════════════════════════════
#  STAT BOX  (large number + label + accent bottom bar)
# ══════════════════════════════════════════════════════════════════
def stat_box(sl, l, t, w, h, num, label, accent=BLUE):
    R(sl, l, t, w, h, fill=_dim(accent, 8), line=_dim(accent, 3), lw=0.6, radius=True)
    R(sl, l, t+h-0.09, w, 0.09, fill=_dim(accent, 2))    # bottom bar
    T(sl, num,   l, t+0.12, w, 0.52, size=30, bold=True, color=accent, align=PP_ALIGN.CENTER)
    T(sl, label, l, t+0.66, w, 0.42, size=9.0, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  INFO BOX  (callout with accent border)
# ══════════════════════════════════════════════════════════════════
def info_box(sl, l, t, w, h, title, body, accent=BLUE):
    R(sl, l, t, w, h, fill=_dim(accent, 9), line=_dim(accent, 3), lw=0.6, radius=True)
    T(sl, title, l+0.15, t+0.10, w-0.28, 0.28, size=11, bold=True, color=accent)
    T(sl, body,  l+0.15, t+0.42, w-0.28, max(0.25, h-0.52), size=9.5, color=MUTED)

# ══════════════════════════════════════════════════════════════════
#  CODE BLOCK  (terminal: dots header + monospace body)
# ══════════════════════════════════════════════════════════════════
def code_box(sl, lines, l, t, w, h, lang=''):
    R(sl, l, t, w, h, fill=CODBG, line=CODBDR, lw=0.6)
    R(sl, l, t, w, 0.30, fill=CODBAR)                     # header bar
    hline(sl, l, t+0.30, w, color=CODBDR)
    # Traffic-light dots
    for di, dc in enumerate([RGBColor(0xff,0x5f,0x57),
                              RGBColor(0xff,0xbd,0x2e),
                              RGBColor(0x28,0xca,0x41)]):
        R(sl, l+0.12+di*0.19, t+0.10, 0.11, 0.11, fill=dc, radius=True)
    if lang:
        T(sl, lang, l+0.75, t+0.05, w-0.9, 0.22,
          size=8, color=RGBColor(0x58,0x6e,0x8a), align=PP_ALIGN.LEFT)
    T(sl, '\n'.join(lines), l+0.14, t+0.36, w-0.26, h-0.42,
      size=8.5, color=CODTXT)

# ══════════════════════════════════════════════════════════════════
#  FLOW NODE  (optionally numbered circle on left)
# ══════════════════════════════════════════════════════════════════
def flow_node(sl, text, l, t, w=2.0, h=0.44, accent=BLUE, num=None):
    R(sl, l, t, w, h, fill=_dim(accent, 7), line=accent, lw=0.6, radius=True)
    if num is not None:
        R(sl, l, t, h, h, fill=accent, radius=True)
        T(sl, str(num), l, t, h, h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        T(sl, text, l+h+0.1, t+0.04, w-h-0.14, h-0.08,
          size=9.5, color=WHITE)
    else:
        T(sl, text, l+0.1, t+0.04, w-0.18, h-0.08,
          size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

def flow_arrow_h(sl, l, t, h=0.44):
    T(sl, '→', l, t, 0.30, h, size=13, color=MUTED, align=PP_ALIGN.CENTER)

def flow_arrow_v(sl, l, t):
    T(sl, '↓', l, t, 0.38, 0.28, size=13, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
#  TABLE
# ══════════════════════════════════════════════════════════════════
def table(sl, headers, rows, l, t, w, h, accent=BLUE, col_widths=None):
    ncols = len(headers); nrows = len(rows) + 1
    tbl = sl.shapes.add_table(
        nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    else:
        cw = w / ncols
        for i in range(ncols): tbl.columns[i].width = Inches(cw)
    # Header row
    for ci, ht in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = _dim(accent, 3)
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = ht
        r.font.size = Pt(9.5); r.font.bold = True
        r.font.color.rgb = accent; r.font.name = 'Calibri'
    # Body rows – zebra
    for ri, row in enumerate(rows):
        rbg = CARD if ri % 2 == 0 else BG3
        for ci, ct in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = rbg
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(ct)
            r.font.size = Pt(10); r.font.bold = (ci == 0)
            r.font.color.rgb = WHITE if ci == 0 else MUTED
            r.font.name = 'Calibri'

# ══════════════════════════════════════════════════════════════════
#  BULLETS  (▸ title + indented description)
# ══════════════════════════════════════════════════════════════════
def bullets(sl, items, l, t, w, h, accent=BLUE, size=12):
    bx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    bx.word_wrap = True; tf = bx.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(4)
        if isinstance(item, tuple):
            title, desc = item
            r = p.add_run(); r.text = '▸  '
            r.font.size = Pt(size-1); r.font.color.rgb = accent; r.font.name = 'Calibri'
            r2 = p.add_run(); r2.text = title
            r2.font.size = Pt(size); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = 'Calibri'
            if desc:
                p2 = tf.add_paragraph(); p2.space_before = Pt(1)
                r3 = p2.add_run(); r3.text = '      ' + desc
                r3.font.size = Pt(size-1.5); r3.font.color.rgb = MUTED; r3.font.name = 'Calibri'
        else:
            r = p.add_run(); r.text = '▸  ' + item
            r.font.size = Pt(size); r.font.color.rgb = MUTED; r.font.name = 'Calibri'

# ══════════════════════════════════════════════════════════════════
#  PROGRESS ROW
# ══════════════════════════════════════════════════════════════════
def progress_row(sl, label, pct, l, t, w, accent=BLUE):
    T(sl, label, l, t, 3.5, 0.30, size=10.5, color=LIGHT)
    track_l = l + 3.65; track_w = w - 4.3
    R(sl, track_l, t+0.09, track_w, 0.12, fill=STROKE)
    fill_w = track_w * pct / 100
    if fill_w > 0.04:
        R(sl, track_l, t+0.09, fill_w, 0.12, fill=accent)
    T(sl, f'{pct}%', l+w-0.55, t, 0.52, 0.30, size=9.5, color=accent, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════
#  CHECK ROW  (✓ or ⚠ with label)
# ══════════════════════════════════════════════════════════════════
def check_row(sl, text, l, t, w, passed=True):
    icon = '✓' if passed else '!'
    acc  = GREEN if passed else YELLOW
    R(sl, l, t, 0.30, 0.30, fill=_dim(acc, 5), line=acc, lw=0.5, radius=True)
    T(sl, icon, l, t, 0.30, 0.30, size=9.5, bold=True, color=acc, align=PP_ALIGN.CENTER)
    T(sl, text, l+0.38, t+0.02, w-0.44, 0.28, size=10.5, color=MUTED)

# ══════════════════════════════════════════════════════════════════
#  SECTION DIVIDER  (two-tone bg + large watermark + title)
# ══════════════════════════════════════════════════════════════════
def section_slide(prs, n, title, subtitle, accent=PURPLE, sn=1):
    sl = new_slide(prs, RGBColor(0x06, 0x08, 0x12))
    R(sl, 0, 0, 7.5, 7.5, fill=_dim(accent, 9))           # tinted left panel
    R(sl, 11.8, 0, 1.533, 7.5, fill=_dim(accent, 7))       # right accent strip
    # Giant watermark number
    T(sl, str(n), 7.5, 0.2, 4.5, 7.1, size=175, bold=True,
      color=RGBColor(0x0d, 0x15, 0x28), align=PP_ALIGN.CENTER)
    badge(sl, f'Section  {n}', 0.65, 1.55, accent)
    T(sl, title, 0.65, 2.05, 10.9, 2.1, size=52, bold=True, color=WHITE)
    R(sl, 0.65, 4.22, 4.2, 0.08, fill=accent)              # accent underline
    T(sl, subtitle, 0.65, 4.42, 10.9, 0.85, size=14.5, color=MUTED)
    snum(sl, sn)
    return sl

# ══════════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── SLIDE 01  TITLE ────────────────────────────────────────────────
sl = new_slide(prs)
R(sl, 0, 0, 7.6, 7.5, fill=RGBColor(0x07, 0x0e, 0x24))   # left pane
R(sl, 7.6, 0, 5.733, 7.5, fill=RGBColor(0x09, 0x06, 0x1e)) # right pane
R(sl, 0, 0, 0.18, 7.5, fill=BLUE)                          # left edge accent
R(sl, 7.58, 0, 0.04, 7.5, fill=STROKE)                     # divider
badge(sl, 'Conference Presentation  2026', 0.55, 0.55, BLUE)
T(sl, 'CIRA', 0.55, 1.05, 7.0, 1.65, size=88, bold=True, color=WHITE)
T(sl, 'Intelligent Conversational AI',
  0.55, 2.7, 7.0, 0.82, size=34, bold=True, color=BLUE)
R(sl, 0.55, 3.58, 3.5, 0.07, fill=BLUE)
T(sl, 'A production-ready, personalized chatbot platform\npowered by Google Gemma, RAG knowledge retrieval\nand deep user profiling — built for scale.',
  0.55, 3.78, 7.0, 1.1, size=13.5, color=MUTED)
tags = [('Google Gemma API',BLUE),('RAG / BM25',PURPLE),
        ('React 18',CYAN),('Node.js + Express',GREEN),
        ('ElevenLabs TTS',YELLOW),('Google Fit',RED)]
x = 0.55
for tag, col in tags:
    badge(sl, tag, x, 5.25, col); x += len(tag)*0.107 + 0.72
# Right side: quick stats
for i, (num, lbl, col) in enumerate([
        ('25+','REST Endpoints',BLUE), ('6','AI Models',PURPLE),
        ('0','Vector DBs',GREEN),      ('~$12','Monthly Cost',YELLOW)]):
    stat_box(sl, 8.0+i%2*2.55, 1.1+i//2*2.0, 2.25, 1.75, num, lbl, col)
snum(sl, 1)

# ── SLIDE 02  AGENDA ───────────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Agenda', 'What We Will Cover Today', PURPLE)
items = [
    ('1','Introduction to CIRA','Overview, mission, tech stack, use cases',BLUE),
    ('2','Core Features','Chat, RAG, TTS, Admin, Multi-instance',PURPLE),
    ('3','Personalization Engine','Profiles, AI Memory, Google Fit',CYAN),
    ('4','Flow Charts','Auth, Chat, RAG, Memory, OAuth flows',GREEN),
    ('5','System Design','DB schema, APIs, JWT, WAL, deployment',YELLOW),
    ('6','Architecture & Advantages','Design patterns, trade-offs, roadmap',RED),
]
for i, (n, title, sub, accent) in enumerate(items):
    c = i % 2; rr = i // 2
    l = 0.48 + c*6.48; t = 1.52 + rr*1.68
    R(sl, l, t, 6.15, 1.52, fill=_dim(accent, 9), line=_dim(accent, 4), lw=0.5, radius=True)
    R(sl, l, t, 0.6,  1.52, fill=_dim(accent, 4), radius=True)  # num background
    T(sl, n, l, t, 0.6, 1.52, size=24, bold=True, color=accent, align=PP_ALIGN.CENTER)
    T(sl, title, l+0.72, t+0.16, 5.25, 0.36, size=13.5, bold=True, color=WHITE)
    T(sl, sub,   l+0.72, t+0.56, 5.25, 0.82, size=10.5, color=MUTED)
T(sl, '70 slides  ·  6 sections  ·  Full Q&A at the end',
  0.48, 6.92, 12.37, 0.38, size=10, color=DIM, align=PP_ALIGN.CENTER)
snum(sl, 2)

# ── SLIDE 03  PROBLEM STATEMENT ────────────────────────────────────
sl = new_slide(prs, RGBColor(0x10, 0x06, 0x0b))
slide_header(sl, 'Problem Statement', 'The Gap in AI Chatbot Solutions', RED)
cards03 = [
    ('🌐','Generic Responses','Most chatbots give one-size-fits-all answers with no knowledge of who the user is or their personal context.'),
    ('🧱','Siloed Knowledge','AI models cannot access private organizational documents without fine-tuning — which is costly and slow to deploy.'),
    ('🔒','Data Hallucination','Large language models frequently generate plausible-sounding but incorrect facts outside their training data.'),
    ('💸','High Infrastructure Cost','Vector databases, embedding APIs, and managed AI services drive up operational costs for small teams.'),
]
for i, (icon, title, desc) in enumerate(cards03):
    c = i % 2; rr = i // 2
    card(sl, 0.48+c*6.48, 1.52+rr*2.68, 6.15, 2.52, icon, title, desc, RED)
R(sl, 0.48, 6.6, 12.37, 0.68, fill=_dim(BLUE,9), line=_dim(BLUE,4), lw=0.5, radius=True)
T(sl, '✦  CIRA solves all four with a single, integrated, open-source architecture.',
  0.7, 6.7, 12.0, 0.48, size=12.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
snum(sl, 3)

# ── SLIDE 04  WHAT IS CIRA? ────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Introduction', 'What is CIRA?', BLUE)
T(sl, 'Conversational Intelligent Responsive Assistant — a full-stack AI chatbot for personalized, knowledge-grounded conversations.',
  0.48, 1.48, 12.37, 0.40, size=12, color=MUTED)
layers04 = [
    ('🤖','AI Core',     'Google Gemma API with 5 model variants (Flash, Pro, open-weight). Configurable temperature presets. Dynamic system prompt injection. Web search grounding. Zero code changes to switch models.',BLUE),
    ('📚','Knowledge',   'BM25 Retrieval-Augmented Generation on private documents (PDF, TXT, Markdown). 400-word chunks, 60-word overlap, 47-stopword filter. Zero embedding API calls. K=5 top chunks per query.',PURPLE),
    ('👤','Personalization','User health profiles (10 fields), auto-extracted AI memories (Gemma analyses 150 messages → 12 facts), and Google Fit fitness snapshot (steps, HR, weight) all injected per-request.',GREEN),
    ('🛠','Admin Control','Real-time bot configuration: model, temperature, system prompt, knowledge weights, TTS settings, KB management, user audit — all via UI, zero redeploy required.',YELLOW),
    ('🔊','Rich UX',     '11-voice ElevenLabs TTS + browser-native STT + full GitHub-Flavored Markdown rendering + dark/light theme switching + fully responsive mobile-first CSS design.',CYAN),
    ('🏗','Multi-Instance','Multiple isolated chatbot instances from one codebase: separate SQLite databases, unique TCP ports, independent data directories, distinct configurations per instance.',RED),
]
for i, (icon, title, desc, acc) in enumerate(layers04):
    c = i % 2; rr = i // 2
    arch_box(sl, 0.48+c*6.48, 1.98+rr*1.78, 6.15, icon, title, desc, acc, h=1.65)
snum(sl, 4)

# ── SLIDE 05  MISSION & VISION ─────────────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'Mission & Vision', 'Our Mission and Vision', CYAN)
R(sl, 0.48, 1.52, 12.37, 2.4, fill=_dim(CYAN,8), line=_dim(CYAN,3), lw=0.6, radius=True)
R(sl, 0.48, 1.52, 0.08, 2.4, fill=CYAN)
T(sl, '🎯  Mission', 0.7, 1.62, 12.0, 0.38, size=16, bold=True, color=CYAN)
T(sl, 'To make intelligent, personalized AI assistants accessible to any organization — without requiring vector infrastructure, GPU clusters, or ML expertise.',
  0.7, 2.06, 11.8, 1.1, size=14, color=LIGHT)
R(sl, 0.48, 4.12, 12.37, 2.4, fill=_dim(PURPLE,8), line=_dim(PURPLE,3), lw=0.6, radius=True)
R(sl, 0.48, 4.12, 0.08, 2.4, fill=PURPLE)
T(sl, '🔭  Vision', 0.7, 4.22, 12.0, 0.38, size=16, bold=True, color=PURPLE)
T(sl, 'A world where every business can deploy a context-aware AI assistant that truly knows its users, grounded in domain knowledge, deployable in minutes on a $7/month server.',
  0.7, 4.66, 11.8, 1.1, size=14, color=LIGHT)
snum(sl, 5)

# ── SLIDE 06  KEY STATS ────────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Key Highlights', 'CIRA by the Numbers', BLUE)
stats06 = [
    ('25+','REST API Endpoints',BLUE),    ('6','Gemma Model Options',PURPLE),
    ('11','ElevenLabs Voices',CYAN),      ('0','Vector DBs Required',GREEN),
    ('7-day','JWT Token Lifetime',YELLOW),('12','Auto-Extracted Memories',RED),
    ('400w','RAG Chunk Size',BLUE),       ('20 MB','Max Doc Upload',PURPLE),
    ('∞','Multi-Instance Support',CYAN),
]
for i, (num, lbl, col) in enumerate(stats06):
    c = i % 3; rr = i // 3
    stat_box(sl, 0.48+c*4.3, 1.52+rr*2.18, 4.05, 1.98, num, lbl, col)
snum(sl, 6)

# ── SLIDE 07  TECH STACK ───────────────────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'Technology Stack', 'Full Technology Stack', CYAN)
T(sl, 'BACKEND', 0.48, 1.52, 5.9, 0.30, size=9.5, bold=True, color=BLUE)
table(sl, ['Component','Technology / Version'],
      [['Runtime','Node.js 20.x LTS'],
       ['Framework','Express.js 4.19.2'],
       ['AI Model','Google Gemma API (5 model variants)'],
       ['RAG Engine','BM25 — pure JS, no external API'],
       ['Database','SQLite 3 + better-sqlite3 9.x (WAL)'],
       ['Auth','JWT (jsonwebtoken) + bcryptjs 2.x'],
       ['File Upload','multer 1.4 — PDF/TXT/MD, 20 MB cap'],
       ['TTS','ElevenLabs API — 11 voices, 4 model tiers'],
       ['Fitness','Google Fit REST API (OAuth2 read-only)']],
      0.48, 1.86, 5.9, 3.60, BLUE)
T(sl, 'FRONTEND', 7.1, 1.52, 5.8, 0.30, size=9.5, bold=True, color=PURPLE)
table(sl, ['Component','Technology / Version'],
      [['Framework','React 18.3.1'],
       ['Routing','React Router v7 (SPA hash routing)'],
       ['Styling','CSS Custom Properties (design tokens)'],
       ['Markdown','react-markdown 9.0.1 + remark-gfm'],
       ['State','React Context API (AuthContext only)'],
       ['Build','Create React App 5.0.1'],
       ['Speech Input','Web Speech API (browser-native)'],
       ['Deploy','Netlify CDN (frontend) / Render.com (API)']],
      7.1, 1.86, 5.8, 3.60, PURPLE)
info_box(sl, 0.48, 5.22, 12.37, 0.72,
         '🏆  Why This Stack?',
         'Zero operational overhead: no vector DB, no ML server, no GPU. Single npm install → deploy in minutes. SQLite scales to ~50 concurrent users with WAL mode.',
         CYAN)
snum(sl, 7)

# ── SLIDE 08  USE CASES ────────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Use Cases', 'Who Uses CIRA?', PURPLE)
cases08 = [
    ('🏥','Healthcare Assistants','Personalized wellness guidance using user medical history, medications, allergies, and live Google Fit fitness data.',BLUE),
    ('🏢','Enterprise Support Bots','Ground AI responses in company documents (PDFs, policies) — accurate, no hallucinations, fully auditable.',PURPLE),
    ('🎓','Education Tutors','Upload course materials; the chatbot answers questions grounded in those specific documents for accurate tutoring.',CYAN),
    ('🏋️','Fitness Coaches','Combine user profile (goals, weight, age) with real-time step and heart rate data for truly personalized coaching.',GREEN),
    ('🏦','Financial Advisors','RAG on financial documents + user profile data for compliant, personalized financial Q&A with source citations.',YELLOW),
    ('⚙️','Internal Knowledge Bases','Replace generic chatbots with a bot that knows your internal docs and processes, with exact source references.',RED),
]
for i, (icon, title, desc, acc) in enumerate(cases08):
    c = i % 3; rr = i // 3
    card(sl, 0.48+c*4.3, 1.52+rr*2.72, 4.05, 2.55, icon, title, desc, acc)
snum(sl, 8)

# ── SLIDE 09  DEPLOYMENT ───────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Deployment', 'Deployment Architecture', BLUE)
arch_box(sl, 0.48, 1.52, 6.0, '☁️', 'Render.com  (Primary)', 'Node.js runtime · 1 GB persistent disk · Auto-deploy from GitHub · Health check endpoint · Env variable management', BLUE, h=0.94)
arch_box(sl, 0.48, 2.56, 6.0, '🌐', 'Netlify  (Frontend Mirror)', 'Static React build · Global CDN · Instant cache invalidation · netlify.toml SPA redirect rules for React Router', PURPLE, h=0.94)
arch_box(sl, 0.48, 3.60, 6.0, '💻', 'Local / Self-Hosted', '`node start-instances.js` launches all instances simultaneously · Isolated data directories per instance · Dev-ready', GREEN, h=0.94)
code_box(sl, [
    '# Required Environment Variables',
    'GEMINI_API_KEY=sk-...',
    'ELEVENLABS_API_KEY=sk_...',
    'JWT_SECRET=long-random-32-char-string',
    '',
    'PORT=5000',
    'NODE_ENV=production',
    'INSTANCE_NAME="CIRA"',
    '',
    '# Google Fit OAuth2',
    'GOOGLE_CLIENT_ID=...',
    'GOOGLE_CLIENT_SECRET=...',
    'GOOGLE_REDIRECT_URI=https://...'],
    6.72, 1.52, 6.13, 3.35, lang='.env')
snum(sl, 9)

# ── SLIDE 10  SECTION 2 DIVIDER ────────────────────────────────────
section_slide(prs, 2, 'Core Features',
    'Multi-conversation chat  ·  RAG knowledge base  ·  Speech I/O  ·  Admin dashboard  ·  Multi-instance',
    PURPLE, 10)

# ── SLIDE 11  FEATURES OVERVIEW ───────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Features Overview', 'Six Pillars of CIRA Functionality', PURPLE)
feats11 = [
    ('💬','Conversation Engine','Multi-session management, auto-titling, persistent history, 10-message rolling context window.',BLUE),
    ('📖','RAG Knowledge Base','BM25-powered keyword retrieval from uploaded PDF, TXT, Markdown. Zero embedding API calls.',PURPLE),
    ('🔊','Speech I/O','Browser-native STT input + ElevenLabs high-fidelity TTS with 11 voice choices and 4 model tiers.',CYAN),
    ('🛡️','Admin Dashboard','Real-time bot settings, user management, KB uploads, conversation audit — full admin control.',GREEN),
    ('👤','Personalization','User profiles + AI memory extraction + Google Fit for truly contextual responses.',YELLOW),
    ('🏗️','Multi-Instance','Multiple isolated chatbot instances from one codebase with separate configs, ports and data.',RED),
]
for i, (icon, title, desc, acc) in enumerate(feats11):
    c = i % 3; rr = i // 3
    card(sl, 0.48+c*4.3, 1.52+rr*2.72, 4.05, 2.55, icon, title, desc, acc)
snum(sl, 11)

# ── SLIDE 12  CONVERSATION MANAGEMENT ─────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'Feature', 'Multi-Conversation Management', CYAN)
bullets(sl, [
    ('Unlimited independent threads', 'Users create, switch between, and manage as many conversations as needed — no per-user caps.'),
    ('Auto-generated titles',         'First user message triggers POST /api/title → Gemma generates a concise 4–6 word title; stored immediately.'),
    ('Persistent SQLite storage',     'All messages stored with Unix-millisecond timestamps and role (user/assistant) in the messages table — fully durable.'),
    ('10-message rolling context',    'Last 10 messages (5 turns) passed as conversation history to Gemma — balances context depth vs. token cost.'),
    ('UUID-based IDs',                'All conversation and message IDs use crypto.randomUUID() — globally unique with zero collision risk.'),
    ('Cascade deletes',               'Deleting a conversation removes all its messages via SQLite ON DELETE CASCADE — no orphan rows ever accumulate.'),
    ('Sidebar navigation',            'Collapsible conversation list with live search. Sidebar auto-collapses to icon-only on screens < 768 px.'),
    ('Rename / delete support',       'PATCH /api/conversations/:id renames, DELETE removes. Both operations immediately reflected across all connected clients.'),
    ('Full admin audit trail',        'Admin can view any user\'s complete conversation history via GET /api/admin/conversations with user filter params.'),
], 0.48, 1.52, 5.9, 5.68, CYAN)
code_box(sl, [
    'CREATE TABLE conversations (',
    '  id         TEXT PRIMARY KEY,',
    '  user_id    TEXT NOT NULL',
    '             REFERENCES users(id)',
    '             ON DELETE CASCADE,',
    "  title      TEXT DEFAULT 'New Chat',",
    '  created_at INTEGER NOT NULL,',
    '  updated_at INTEGER NOT NULL',
    ');',
    '',
    'CREATE TABLE messages (',
    '  id              TEXT PRIMARY KEY,',
    '  conversation_id TEXT NOT NULL',
    '    REFERENCES conversations(id)',
    '    ON DELETE CASCADE,',
    '  role      TEXT,  -- user | assistant',
    '  text      TEXT NOT NULL,',
    '  timestamp INTEGER NOT NULL,',
    '  refs      TEXT  -- JSON KB references',
    ');'],
    6.72, 1.52, 6.13, 5.68, lang='schema.sql')
snum(sl, 12)

# ── SLIDE 13  AI RESPONSE ENGINE ──────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'AI Engine', 'AI Response Engine', BLUE)
arch_box(sl, 0.48, 1.52, 5.9, '🤖', 'Gemma Model Flexibility', 'Supports gemma-3-flash-preview, gemma-2.5-pro, gemma-2.5-flash, gemma-2.0-flash, and gemma-4 open-weight variant. Admin switches model live from the settings dropdown — no code change or server restart required.', BLUE, h=0.96)
arch_box(sl, 0.48, 2.58, 5.9, '🌡️', 'Temperature Control', 'Three configurable presets: Precise (0.2) for factual/compliance use, Balanced (0.7) for everyday chat, Creative (1.2) for brainstorming. Stored in the settings table and applied per-request.', PURPLE, h=0.96)
arch_box(sl, 0.48, 3.64, 5.9, '📝', 'System Prompt Injection', 'A dynamic systemInstruction is built for every request: base admin prompt + user profile + extracted memories + fitness snapshot + Top-K RAG chunks + last 10 messages. Total assembly time < 20 ms.', CYAN, h=0.96)
arch_box(sl, 0.48, 4.70, 5.9, '🌐', 'Web Search Tool', "Gemma API's built-in Google Search grounding provides live web results. Weight configurable 0–100% via admin slider. Automatically disabled when Gemma open-weight models are selected (they lack tool-use support).", GREEN, h=0.96)
T(sl, 'Knowledge Source Weights  (0–100%, Admin Configurable)', 6.72, 1.52, 6.13, 0.34, size=11, bold=True, color=BLUE)
for i, (lbl, pct, col) in enumerate([
        ('KB Documents (RAG)',   80, BLUE),
        ("Gemma's Own Knowledge", 70, PURPLE),
        ('Live Web Search',      40, CYAN)]):
    progress_row(sl, lbl, pct, 6.72, 1.98+i*0.54, 6.13, col)
info_box(sl, 6.72, 3.72, 6.13, 0.84,
         '⚠  Gemma Open-Weight Limitation',
         'Gemma 4 and other open-weight variants do not support the Google Search grounding tool. webSearchWeight > 0 with Gemma is silently ignored at runtime. The admin settings UI displays a warning badge when a Gemma model is selected.', YELLOW)
info_box(sl, 6.72, 4.68, 6.13, 0.84,
         '✅  Zero-Restart Configuration',
         'All model, temperature, knowledge weight, and system prompt changes are read from the SQLite settings table at the start of every chat request — changes take effect immediately with no server restart or redeploy.', GREEN)
snum(sl, 13)

# ── SLIDE 14  RAG KNOWLEDGE BASE ──────────────────────────────────
sl = new_slide(prs, RGBColor(0x06, 0x12, 0x0c))
slide_header(sl, 'RAG System', 'Knowledge Base & RAG Pipeline', GREEN)
bullets(sl, [
    ('Supported formats',   'PDF (text extraction via pdf-parse 3.1.0), plain TXT, and Markdown files up to 20 MB per upload'),
    ('Admin upload UI',     'Drag-and-drop KnowledgePage with live progress, chunk count display, and delete-by-document control'),
    ('Smart chunking',      '400-word segments with 60-word overlap — overlap preserves sentence context across chunk boundaries'),
    ('BM25 tokenization',   'Tokens stored as JSON string arrays in kb_chunks.embedding column — zero external embedding API calls'),
    ('47 stopwords filtered','Common words (the, is, at, which…) removed before indexing to improve retrieval signal-to-noise ratio'),
    ('Top-K retrieval',     'Default K=5 highest-scoring BM25 chunks per query; admin-configurable via ragTopK setting'),
    ('Source citations',    'Document name + chunk preview injected alongside each chunk; shown as references in AI response UI'),
    ('Text snippets',       'Admin can also add free-text snippets (FAQ items, policy statements) directly — no file upload needed'),
    ('Configurable weight', 'RAG influence 0–100% independently tunable in admin settings, separate from own-knowledge and web-search weights'),
    ('KB-only strict mode', 'Set ragWeight=100, ownKnowledgeWeight=0, webSearchWeight=0 for fully grounded, hallucination-free answers'),
], 0.48, 1.52, 5.9, 5.68, GREEN)
T(sl, 'Ingestion & Retrieval Pipeline', 6.72, 1.52, 6.13, 0.32, size=11, bold=True, color=GREEN)
pipe14 = [
    ('Upload PDF / TXT / MD', GREEN),
    ('Extract raw text content', CYAN),
    ('Chunk: 400w, 60w overlap', BLUE),
    ('BM25 tokenize → store JSON', PURPLE),
    ('Query → BM25 score all chunks', YELLOW),
    ('Return Top-K + source names', GREEN),
    ('Inject into Gemma systemInstruction', CYAN),
]
for i, (txt, acc) in enumerate(pipe14):
    flow_node(sl, txt, 6.72, 1.92+i*0.76, 6.13, 0.52, acc, num=i+1)
    if i < 6: flow_arrow_v(sl, 9.58, 2.44+i*0.76)
snum(sl, 14)

# ── SLIDE 15  BM25 ALGORITHM ──────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'BM25 Algorithm', 'BM25 — Why No Vector Embeddings?', BLUE)
R(sl, 0.48, 1.52, 5.9, 1.28, fill=_dim(BLUE,9), line=_dim(BLUE,3), lw=0.5, radius=True)
R(sl, 0.48, 1.52, 0.08, 1.28, fill=BLUE)
T(sl, 'BM25 Score Formula', 0.68, 1.60, 5.6, 0.30, size=11, bold=True, color=BLUE)
T(sl, 'score(D,Q) = Σ IDF(qi) × [ tf(qi,D) × (k1+1) ] / [ tf(qi,D) + k1 × (1−b+b×|D|/avgdl) ]',
  0.68, 1.96, 5.6, 0.72, size=10, color=LIGHT)
T(sl, 'K1=1.5 (term saturation)  ·  B=0.75 (length norm)  ·  47 stopwords  ·  Top-K=5 returned  ·  IDF computed over full corpus',
  0.68, 2.60, 5.6, 0.32, size=8.5, color=MUTED)
table(sl, ['Aspect','BM25 (CIRA)','Vector Embeddings'],
      [['API calls','None','Per chunk'],['Storage','~12 MB SQLite','GBs (FAISS/Pinecone)'],
       ['Search latency','< 5 ms','50–200 ms'],['Monthly cost','$0','$70–250'],
       ['Explainability','High (keywords)','Low (opaque)'],
       ['Semantic search','Keyword-based','Semantic']],
      0.48, 3.0, 5.9, 3.2, BLUE)
code_box(sl, [
    '// embed.js — BM25 constants',
    'const K1 = 1.5;',
    'const B  = 0.75;',
    '',
    'const STOPWORDS = new Set([',
    "  'the','is','at','which','on',",
    "  'a','an','and','or','but',",
    '  // ... 42 more',
    ']);',
    '',
    '// Tokenize → lowercase → filter',
    '// → store as JSON token array',
    '// e.g. ["diabetes","management"]'],
    6.72, 1.52, 6.13, 4.65, lang='embed.js')
info_box(sl, 6.72, 6.3, 6.13, 0.65,
         '✅  Practical Result',
         'For domain-specific documents, BM25 matches semantic search accuracy at zero embedding cost.',
         GREEN)
snum(sl, 15)

# ── SLIDE 16  SPEECH FEATURES ─────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Speech I/O', 'Speech-to-Text & Text-to-Speech', PURPLE)
T(sl, '🎤  Speech-to-Text (STT)', 0.48, 1.52, 5.9, 0.32, size=13, bold=True, color=BLUE)
bullets(sl, [
    'Uses browser-native Web Speech API — zero backend cost',
    'Continuous recognition with live interim results displayed',
    'Graceful fallback when browser does not support it',
    'Works on Chrome, Edge, Safari (WebKit)',
], 0.48, 1.90, 5.9, 1.78, BLUE)
T(sl, '🔊  Text-to-Speech (TTS) — ElevenLabs', 0.48, 3.75, 5.9, 0.32, size=13, bold=True, color=PURPLE)
bullets(sl, [
    ('11+ voice personalities',    'Rachel, Adam, and more studio-quality options'),
    ('4 model tiers',              'Turbo v2.5 (default), Turbo v2, Multilingual v2, Monolingual v1'),
    ('Stability & similarity boost','Fine-tune voice characteristics per deployment need'),
    ('Markdown stripped',          'Clean audio — no asterisks or code fences read aloud'),
    ('MP3 output',                 '44.1 kHz, 128 kbps audio returned by /api/tts endpoint'),
], 0.48, 4.14, 5.9, 3.06, PURPLE)
code_box(sl, [
    "// POST /api/tts  —  server.js",
    "app.post('/api/tts', requireAuth, async (req, res) => {",
    "  const { text, voiceId, modelId,",
    "          stability, similarityBoost } = req.body;",
    "",
    "  // Strip markdown for clean audio",
    "  const clean = stripMarkdown(text);",
    "",
    "  const audio = await fetchElevenLabs({",
    "    text:       clean,",
    "    voice_id:   voiceId,",
    "    model_id:   modelId,",
    "    voice_settings: {",
    "      stability,",
    "      similarity_boost: similarityBoost",
    "    }",
    "  });",
    "  res.set('Content-Type', 'audio/mpeg');",
    "  res.send(audio);",
    "});"],
    6.72, 1.52, 6.13, 5.68, lang='server.js')
snum(sl, 16)

# ── SLIDE 17  ADMIN DASHBOARD ─────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Admin', 'Admin Dashboard & Control Panel', YELLOW)
panels17 = [
    ('📊','Dashboard Analytics','Total users · KB documents · Avg login count · Most recent activity · Active session tracking',YELLOW),
    ('👥','User Management','View all users · Add admin notes · See login history · Grant or revoke admin role',BLUE),
    ('🗂️','Chat History Audit','Full conversation visibility for any user · Compliance and moderation support',PURPLE),
    ('⚙️','Bot Settings','Real-time model selection · System prompt · Temperature · RAG weights · TTS defaults',CYAN),
    ('📚','Knowledge Base','Drag-and-drop upload · Text snippets · View all docs · Delete documents · Chunk counts',GREEN),
    ('🔒','Access Control','All admin routes protected with requireAdmin middleware. JWT role claim verified on every request.',RED),
]
for i, (icon, title, desc, acc) in enumerate(panels17):
    c = i % 3; rr = i // 3
    card(sl, 0.48+c*4.3, 1.52+rr*2.68, 4.05, 2.52, icon, title, desc, acc)
snum(sl, 17)

# ── SLIDE 18  MULTI-INSTANCE ──────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Multi-Instance', 'Multi-Instance Architecture', CYAN)
T(sl, 'A single CIRA codebase powers multiple independent chatbot instances simultaneously, each fully isolated:',
  0.48, 1.52, 5.9, 0.52, size=12, color=MUTED)
bullets(sl, [
    ('Isolated SQLite database', 'Each instance has its own users, conversations, and knowledge base'),
    ('Unique TCP port',          '5000, 5001, 5002 … configurable in instances.json'),
    ('Separate data directory',  'Independent persistent disk location per instance'),
    ('Independent bot name',     'Different INSTANCE_NAME and configuration per instance'),
    ('Single launcher script',   'node start-instances.js spawns all processes at once'),
], 0.48, 2.14, 5.9, 3.72, CYAN)
info_box(sl, 0.48, 5.98, 5.9, 0.72,
         'Use Case',
         'CIRA A for public users, CIRA B for internal staff — same codebase, completely separate databases and configs.',
         CYAN)
code_box(sl, [
    '// instances.json',
    '[',
    '  { "name": "CIRA A",',
    '    "port": 5000,',
    '    "dataDir": "data/chatbot-a" },',
    '  { "name": "CIRA B",',
    '    "port": 5001,',
    '    "dataDir": "data/chatbot-b" }',
    ']',
    '',
    '// start-instances.js',
    'for (const inst of instances) {',
    '  spawn("node", ["server.js"], {',
    '    env: { ...process.env,',
    '      PORT:          inst.port,',
    '      INSTANCE_NAME: inst.name,',
    '      DATA_DIR:      inst.dataDir',
    '    }',
    '  });',
    '}'],
    6.72, 1.52, 6.13, 5.68, lang='instances.json')
snum(sl, 18)

# ── SLIDE 19  UI FEATURES ─────────────────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'UI Features', 'Rich UI & Experience Features', GREEN)
feats19 = [
    ('📝','Markdown Rendering','Full GitHub-Flavored Markdown via react-markdown. Code blocks with syntax highlighting, tables, lists — rendered safely without XSS risk.',GREEN),
    ('🌗','Dark / Light Themes','CSS custom properties (design tokens) enable instant theme switching. Dark mode is the default. Admin sets the default theme for all users.',BLUE),
    ('📱','Fully Responsive','Mobile-first CSS Grid/Flexbox layout. Sidebar collapses on small screens. Touch-friendly controls. Works on phones, tablets, and desktops.',PURPLE),
    ('📎','Source Citations','KB document names appear as clickable references alongside AI responses for full traceability and compliance auditability.',CYAN),
]
for i, (icon, title, desc, acc) in enumerate(feats19):
    c = i % 2; rr = i // 2
    card(sl, 0.48+c*6.48, 1.52+rr*2.68, 6.15, 2.52, icon, title, desc, acc)
snum(sl, 19)

# ── SLIDE 20  SECTION 3 DIVIDER ────────────────────────────────────
section_slide(prs, 3, 'The Personalization Engine',
    'User profiles  ·  AI memory extraction  ·  Google Fit integration  ·  Dynamic context assembly',
    GREEN, 20)

# ── SLIDE 21  PERSONALIZATION PHILOSOPHY ──────────────────────────
sl = new_slide(prs, RGBColor(0x05, 0x10, 0x0b))
slide_header(sl, 'Philosophy', 'Why Personalization Matters', GREEN)
pillars21 = [
    ('🧠','Context-Aware','The AI knows who you are — age, health conditions, goals, and history — before you type a single word.',GREEN),
    ('💾','Memory Persistence','Facts learned in one conversation persist across all future sessions. The bot gets smarter over time automatically.',CYAN),
    ('📡','Live Data','Real-time fitness data from Google Fit — steps, heart rate, weight — refreshed on demand for truly current advice.',BLUE),
]
for i, (icon, title, desc, acc) in enumerate(pillars21):
    card(sl, 0.48+i*4.3, 1.52, 4.05, 3.58, icon, title, desc, acc)
R(sl, 0.48, 5.28, 12.37, 0.98, fill=_dim(GREEN,8), line=_dim(GREEN,3), lw=0.6, radius=True)
R(sl, 0.48, 5.28, 0.08, 0.98, fill=GREEN)
T(sl, 'All personalization data is assembled into a single dynamic systemInstruction sent with every Gemma API call — seamlessly and automatically, in under 20 ms.',
  0.72, 5.40, 11.9, 0.75, size=13, color=LIGHT)
snum(sl, 21)

# ── SLIDE 22  USER PROFILE SYSTEM ─────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'User Profile', 'User Profile System', GREEN)
table(sl, ['Field','Data Type','Example Value / Notes'],
      [['Age','Number (integer)','32 — used for age-appropriate health guidance'],
       ['Gender','String','Female — influences health reference ranges'],
       ['Height','Number (cm)','168 — combined with weight for BMI calculation'],
       ['Weight','Number (kg)','65.0 — compared against Google Fit measurements'],
       ['Medical Conditions','Free text','"Type 2 diabetes" — triggers condition-aware advice'],
       ['Medications','Free text','"Metformin 500 mg" — avoids conflicting recommendations'],
       ['Allergies','Free text','"Peanuts, sulfa drugs" — critical for dietary advice'],
       ['Health Goals','Free text','"Lose 5 kg by June" — used to frame all coaching advice'],
       ['Activity Level','String','Sedentary / Lightly active / Very active'],
       ['Custom Notes','Free text','Any user-specific context the AI should always remember']],
      0.48, 1.52, 6.3, 5.0, GREEN, col_widths=[2.0, 1.55, 2.75])
code_box(sl, [
    '// Profile injected into systemInstruction',
    'function buildSystemPrompt(settings, user) {',
    '  let prompt = settings.systemPrompt;',
    '',
    '  if (user.profile?.age) {',
    '    prompt += `\\nUser Profile:',
    '  Age:         ${user.profile.age}',
    '  Gender:      ${user.profile.gender}',
    '  Conditions:  ${user.profile.conditions}',
    '  Medications: ${user.profile.medications}',
    '  Allergies:   ${user.profile.allergies}',
    '  Goals:       ${user.profile.goals}`;',
    '  }',
    '  return prompt;',
    '}'],
    6.92, 1.52, 5.93, 5.0, lang='server.js')
info_box(sl, 6.92, 6.64, 5.93, 0.62,
         '✅  Gemma Receives Full Context',
         'User data injected automatically — no manual prompting needed.', GREEN)
snum(sl, 22)

# ── SLIDE 23  AI MEMORY SYSTEM ─────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'AI Memory', 'AI Memory Extraction System', PURPLE)
T(sl, 'CIRA automatically learns facts about users by analysing recent conversation history using Gemma:',
  0.48, 1.52, 5.9, 0.48, size=12, color=MUTED)
steps23 = [
    ('Trigger Analysis',   'User triggers POST /api/user/memories/extract. System fetches last 150 messages across ALL conversations for that user.',BLUE),
    ('Gemma Extraction',   'Gemma processes the message history with a structured prompt: "Extract exactly 12 new facts about this user not already in their memory list."',PURPLE),
    ('JSON Parsing',       'Response expected as a JSON array of strings. Markdown code fences (```json...```) are stripped before JSON.parse() to handle common model formatting quirks.',CYAN),
    ('Deduplication',      'Each candidate fact is compared against the first 40 characters of every existing memory — matching prefix = duplicate, skip. Prevents bloat on re-runs.',YELLOW),
    ('Storage & Injection','Unique new facts appended to users.memories (JSON column). On next chat, all memories listed under "What I know about you:" in the systemInstruction.',GREEN),
]
for i, (title, desc, acc) in enumerate(steps23):
    R(sl, 0.48, 2.12+i*1.10, 0.58, 0.58, fill=_dim(acc,4), line=acc, lw=0.6, radius=True)
    T(sl, str(i+1), 0.48, 2.12+i*1.10, 0.58, 0.58, size=14, bold=True, color=acc, align=PP_ALIGN.CENTER)
    R(sl, 1.16, 2.12+i*1.10, 5.22, 0.70, fill=_dim(acc,9), line=_dim(acc,4), lw=0.5, radius=True)
    T(sl, title, 1.30, 2.16+i*1.10, 5.0, 0.26, size=11.5, bold=True, color=WHITE)
    T(sl, desc,  1.30, 2.44+i*1.10, 5.0, 0.42, size=8.8, color=MUTED)
code_box(sl, [
    '// users.memories — JSON array column',
    '[',
    '  {',
    '    "id":        "mem_abc123",',
    '    "text":      "Prefers morning workouts",',
    '    "source":    "auto",  // auto | manual',
    '    "createdAt": 1716890400000',
    '  },',
    '  { "id": "mem_def456",',
    '    "text": "Has lactose intolerance",',
    '    "source": "auto",',
    '    "createdAt": 1716890500000 }',
    ']',
    '',
    '// systemInstruction snippet (injected per chat)',
    '"What I know about you:',
    ' - Prefers morning workouts',
    ' - Has lactose intolerance',
    ' - Training for a 5K run in July',
    ' - Takes vitamin D supplements daily"'],
    6.72, 1.52, 6.13, 5.2, lang='server.js')
info_box(sl, 6.72, 6.28, 6.13, 0.72,
         'Memory Management',
         'View · manually delete · or add custom memories from the Profile Page. Re-run extraction at any time.', PURPLE)
snum(sl, 23)

# ── SLIDE 24  GOOGLE FIT INTEGRATION ──────────────────────────────
sl = new_slide(prs, RGBColor(0x05, 0x12, 0x0c))
slide_header(sl, 'Google Fit', 'Google Fit Integration', GREEN)
fit_items = [
    ('🔐','OAuth2 Authorization','User clicks "Connect Google Fit" → Google consent screen → tokens stored securely server-side per user.',GREEN),
    ('📊','Data Retrieved (7 Days)','Daily step counts · Average heart rate (BPM) · Weight measurements — aggregated over a rolling 7-day window.',CYAN),
    ('🧠','Context Injection','Snapshot appended to Gemma systemInstruction: "7,432 steps/day, 72 BPM avg heart rate, 68.5 kg".',BLUE),
    ('🔄','Refresh on Demand','User manually refreshes fitness data at any time. Cached in user profile with a last-updated timestamp.',YELLOW),
]
for i, (icon, title, desc, acc) in enumerate(fit_items):
    arch_box(sl, 0.48, 1.52+i*1.42, 5.9, icon, title, desc, acc, h=1.30)
code_box(sl, [
    '// Read-only OAuth2 scopes only',
    'const SCOPES = [',
    "  'https://www.googleapis.com/auth/",
    "    fitness.activity.read',",
    "  'https://www.googleapis.com/auth/",
    "    fitness.body.read',",
    "  'https://www.googleapis.com/auth/",
    "    fitness.heart_rate.read'",
    '];',
    '',
    '// Fitness injected into system prompt:',
    '"Recent Fitness Data (last 7 days):',
    ' • Avg daily steps: 7,432',
    ' • Avg heart rate:  72 BPM',
    ' • Latest weight:   68.5 kg',
    'Use this to give relevant advice."'],
    6.72, 1.52, 6.13, 5.68, lang='server.js')
snum(sl, 24)

# ── SLIDE 25  DYNAMIC SYSTEM PROMPT ───────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Context Assembly', 'Dynamic System Prompt Construction', PURPLE)
T(sl, 'Every Gemma API call assembles a rich systemInstruction from 6 data sources in real-time (< 20 ms total):',
  0.48, 1.52, 12.37, 0.40, size=12, color=MUTED)
sources25 = [
    ('⚙️','1. Admin System Prompt','The base instruction block: bot persona, tone, topic scope, and any hardcoded guidelines set by the admin in the settings UI. Read from SQLite on every request.',BLUE),
    ('👤','2. User Health Profile','Age, gender, height, weight, medical conditions, medications, allergies, health goals, activity level, and custom notes — up to 10 structured fields injected as a formatted block.',PURPLE),
    ('💾','3. AI Extracted Memories','Previously learned facts about the user: "What I know about you: - prefers morning workouts - has lactose intolerance…". Extracted by Gemma from chat history, persisted across sessions.',GREEN),
    ('🏃','4. Google Fit Snapshot','7-day rolling fitness data if user has connected Google Fit: average daily steps, average heart rate (BPM), and latest weight measurement — refreshed on demand by the user.',YELLOW),
    ('📚','5. Top-K RAG Chunks','The 5 highest-scoring BM25 chunks from the knowledge base for the current query. Each chunk includes the source document name for citation. Injected with doc names for traceability.',CYAN),
    ('🕐','6. Conversation History','The last 10 messages (5 user turns + 5 assistant turns) from the current conversation, passed as structured history turns to Gemma for multi-turn context and coherence.',RED),
]
for i, (icon, num, desc, acc) in enumerate(sources25):
    c = i % 3; rr = i // 3
    R(sl, 0.48+c*4.3, 2.05+rr*2.42, 4.05, 2.28, fill=_dim(acc,9), line=acc, lw=0.6, radius=True)
    R(sl, 0.48+c*4.3, 2.05+rr*2.42, 0.06, 2.28, fill=acc)
    T(sl, icon+' '+num, 0.65+c*4.3, 2.15+rr*2.42, 3.8, 0.30, size=11, bold=True, color=acc)
    T(sl, desc, 0.65+c*4.3, 2.50+rr*2.42, 3.8, 1.55, size=9.2, color=MUTED)
snum(sl, 25)

# ── SLIDE 26  PERSONALIZATION IMPACT ──────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Impact', 'Personalization in Action — Before vs. After', GREEN)
R(sl, 0.48, 1.52, 5.9, 4.72, fill=_dim(RED,10), line=RED, lw=0.6, radius=True)
R(sl, 0.48, 1.52, 0.08, 4.72, fill=RED)
T(sl, '❌  Generic Chatbot', 0.68, 1.62, 5.6, 0.32, size=12, bold=True, color=RED)
T(sl, 'Q: What should I eat before a workout?', 0.68, 2.02, 5.6, 0.30, size=11, bold=True, color=LIGHT)
T(sl, '"Before a workout, it is generally recommended to eat a balanced meal containing carbohydrates, protein, and healthy fats 2–3 hours before exercising. Examples include oatmeal with fruit, a chicken sandwich, or a banana with peanut butter."',
  0.68, 2.42, 5.6, 2.5, size=10.5, color=MUTED)
R(sl, 6.72, 1.52, 6.13, 4.72, fill=_dim(GREEN,10), line=GREEN, lw=0.6, radius=True)
R(sl, 6.72, 1.52, 0.08, 4.72, fill=GREEN)
T(sl, '✅  CIRA Personalised Response', 6.90, 1.62, 5.9, 0.32, size=12, bold=True, color=GREEN)
T(sl, 'Q: What should I eat before a workout?', 6.90, 2.02, 5.9, 0.30, size=11, bold=True, color=LIGHT)
T(sl, '"Given your Type 2 diabetes and Metformin prescription, avoid simple carbs that spike blood sugar. Greek yogurt with berries works well. Based on your 7,432 daily steps and goal to lose 5 kg, keep pre-workout calories under 250. Since you prefer morning workouts, eat 30–45 minutes before you start."',
  6.90, 2.42, 5.9, 2.5, size=10.5, color=MUTED)
R(sl, 0.48, 6.4, 12.37, 0.72, fill=_dim(BLUE,9), line=BLUE, lw=0.5, radius=True)
T(sl, 'Profile data  +  AI memories  +  fitness snapshot  =  genuinely useful, personalised advice',
  0.7, 6.5, 12.0, 0.52, size=12.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
snum(sl, 26)

# ── SLIDE 27  SECTION 4 DIVIDER ────────────────────────────────────
section_slide(prs, 4, 'Flow Charts',
    'Authentication  ·  Chat processing  ·  RAG pipeline  ·  Memory extraction  ·  OAuth  ·  Error handling',
    YELLOW, 27)

# ── SLIDE 28  AUTH FLOW ────────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Auth Flow', 'User Registration & Login Flow', YELLOW)
T(sl, 'Registration', 0.48, 1.52, 5.9, 0.30, size=11.5, bold=True, color=BLUE)
reg28 = [
    ('User enters email + password + name',BLUE),
    ('Validate: email unique, password ≥ 6 chars',YELLOW),
    ('bcryptjs hash password (10 salt rounds)',PURPLE),
    ('INSERT user row with UUID → SQLite',BLUE),
    ('Sign JWT (7-day expiry) → Return token',GREEN),
    ('Store in localStorage → Route to /chat',CYAN),
]
for i, (txt, acc) in enumerate(reg28):
    flow_node(sl, txt, 0.48, 1.88+i*0.84, 5.9, 0.54, acc, num=i+1)
    if i < 5: flow_arrow_v(sl, 3.23, 2.42+i*0.84)
T(sl, 'Every Protected Request', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=PURPLE)
prot28 = [
    ('Client sends Bearer JWT in Authorization header',BLUE),
    ('requireAuth: jwt.verify(token, JWT_SECRET)',YELLOW),
    ('Admin routes: also check role === "admin"',PURPLE),
]
for i, (txt, acc) in enumerate(prot28):
    flow_node(sl, txt, 6.72, 1.88+i*0.94, 6.13, 0.60, acc, num=i+1)
    if i < 2: flow_arrow_v(sl, 9.59, 2.48+i*0.94)
# Outcome boxes
R(sl, 6.72, 4.72, 2.95, 0.58, fill=_dim(GREEN,8), line=GREEN, lw=0.5, radius=True)
T(sl, '✓  Valid → continue', 6.80, 4.80, 2.8, 0.42, size=10, color=WHITE, align=PP_ALIGN.CENTER)
R(sl, 9.85, 4.72, 3.0, 0.58, fill=_dim(RED,8), line=RED, lw=0.5, radius=True)
T(sl, '✗  Invalid → 401 Unauthorized', 9.93, 4.80, 2.85, 0.42, size=10, color=WHITE, align=PP_ALIGN.CENTER)
snum(sl, 28)

# ── SLIDE 29  CHAT MESSAGE FLOW ────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Chat Flow', 'Chat Message Processing Flow — 12 Steps', CYAN)
steps29 = [
    ('POST /api/chat  {conversationId, message}',  BLUE),
    ('requireAuth — jwt.verify(token, JWT_SECRET)', YELLOW),
    ('SELECT profile + memories FROM users',        PURPLE),
    ('Check fitness snapshot in user.profile',      GREEN),
    ('BM25 search kb_chunks → Top-K chunks',        CYAN),
    ('SELECT last 10 msgs FROM messages ORDER BY timestamp',BLUE),
    ('SELECT all settings FROM settings table',     YELLOW),
    ('Build systemInstruction string (<20 ms total)',PURPLE),
    ('POST Gemma API with system + history',        GREEN),
    ('Receive streamed / complete AI response text',CYAN),
    ('INSERT user msg + assistant msg → SQLite',    BLUE),
    ('Return {text, refs:[{docName, snippet}]}',    GREEN),
]
for i, (txt, acc) in enumerate(steps29):
    c = i % 4; rr = i // 4
    flow_node(sl, txt, 0.42+c*3.24, 1.52+rr*1.96, 3.06, 0.52, acc, num=i+1)
    if c < 3: flow_arrow_h(sl, 3.50+c*3.24, 1.58+rr*1.96, 0.52)
R(sl, 0.48, 7.0, 12.37, 0.30, fill=_dim(BLUE,10), line=DIM, lw=0.4, radius=True)
T(sl, '⚡  Steps 2–8 complete in < 20 ms (SQLite reads). Only the Gemma API call (step 9) adds network latency.',
  0.65, 7.05, 12.0, 0.22, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
snum(sl, 29)

# ── SLIDE 30  RAG PIPELINE FLOW ────────────────────────────────────
sl = new_slide(prs, RGBColor(0x05, 0x12, 0x0c))
slide_header(sl, 'RAG Pipeline', 'RAG Pipeline — Detailed Flow', GREEN)
T(sl, 'Ingestion Phase (Document Upload)', 0.48, 1.52, 5.9, 0.30, size=11.5, bold=True, color=GREEN)
ing30 = [
    ('Admin uploads PDF/TXT/MD file',       GREEN),
    ('Multer validates type & size ≤ 20 MB', YELLOW),
    ('Text extraction (pdf-parse / buffer)', CYAN),
    ('chunker.js: 400w segments, 60w overlap',BLUE),
    ('embed.js: BM25 tokenize each chunk',   PURPLE),
    ('Store JSON tokens in kb_chunks table', GREEN),
]
for i, (txt, acc) in enumerate(ing30):
    flow_node(sl, txt, 0.48, 1.88+i*0.84, 5.9, 0.54, acc, num=i+1)
    if i < 5: flow_arrow_v(sl, 3.23, 2.42+i*0.84)
T(sl, 'Retrieval Phase (Chat Query)', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=CYAN)
ret30 = [
    ('Tokenize user query (same BM25 tokenizer)', CYAN),
    ('Score all KB chunks with BM25 formula',     BLUE),
    ('Sort by relevance score descending',         YELLOW),
    ('Return Top-5 chunks + source doc names',    GREEN),
    ('Inject into Gemma systemInstruction',      PURPLE),
    ('AI response cites source documents',        CYAN),
]
for i, (txt, acc) in enumerate(ret30):
    flow_node(sl, txt, 6.72, 1.88+i*0.84, 6.13, 0.54, acc, num=i+1)
    if i < 5: flow_arrow_v(sl, 9.59, 2.42+i*0.84)
snum(sl, 30)

# ── SLIDE 31  MEMORY EXTRACTION FLOW ──────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Memory Flow', 'AI Memory Extraction Flow', PURPLE)
flow31 = [
    ('User triggers POST /api/user/memories/extract', BLUE),
    ('Authenticate JWT, verify user identity',        YELLOW),
    ('Fetch last 150 messages from all conversations',PURPLE),
    ('Build extraction prompt: "Extract 12 new user facts"', CYAN),
    ('Call Gemma API in JSON mode (structured output)',      GREEN),
    ('Parse JSON (strip markdown fences if present)',         BLUE),
    ('Dedup: 40-char prefix match vs. existing memories',    YELLOW),
    ('Append new facts to users.memories JSON column',        PURPLE),
    ('Return updated memory list to client',                  GREEN),
]
for i, (txt, acc) in enumerate(flow31):
    c = i % 3; rr = i // 3
    flow_node(sl, txt, 0.42+c*4.33, 1.52+rr*1.52, 4.06, 0.72, acc, num=i+1)
    if c < 2: flow_arrow_h(sl, 4.50+c*4.33, 1.70+rr*1.52, 0.72)
info_box(sl, 0.48, 6.12, 12.37, 0.68,
         'Fallback Strategy',
         'If Gemma returns JSON wrapped in markdown fences, they are stripped before JSON.parse(). On parse failure, falls back to Gemma 1.5 Flash in forced JSON mode.',
         YELLOW)
snum(sl, 31)

# ── SLIDE 32  GOOGLE FIT OAUTH FLOW ───────────────────────────────
sl = new_slide(prs, RGBColor(0x05, 0x12, 0x0c))
slide_header(sl, 'OAuth Flow', 'Google Fit OAuth2 Flow', GREEN)
oauth32 = [
    ('User: GET /api/fitness/auth-url',              BLUE),
    ('Server builds OAuth2 URL + state param',       CYAN),
    ('User redirected to Google consent screen',     GREEN),
    ('User grants read-only permission on Google',   YELLOW),
    ('Google redirects to /api/fitness/callback?code=...', BLUE),
    ('Exchange code → access_token + refresh_token', PURPLE),
    ('Tokens stored in user.profile (SQLite)',       GREEN),
    ('POST /api/fitness/refresh → fetch 7-day data', CYAN),
    ('Steps + HR + Weight cached in user profile',   YELLOW),
    ('Fitness snapshot injected into systemInstruction', GREEN),
]
for i, (txt, acc) in enumerate(oauth32):
    c = i % 2; rr = i // 2
    flow_node(sl, txt, 0.48+c*6.57, 1.52+rr*1.04, 6.3, 0.68, acc, num=i+1)
    if c == 0: flow_arrow_h(sl, 6.80, 1.68+rr*1.04, 0.68)
info_box(sl, 0.48, 6.65, 12.37, 0.62,
         '🔒  Security',
         'Only read-only scopes requested. Tokens stored server-side only. Revoke anytime via DELETE /api/fitness/disconnect.',
         GREEN)
snum(sl, 32)

# ── SLIDE 33  DOCUMENT UPLOAD FLOW ────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Document Upload', 'Document Upload & Processing Flow', BLUE)
doc33 = [
    ('Admin: POST /api/documents (multipart)',          BLUE),
    ('Multer: validate file type & size ≤ 20 MB',       YELLOW),
    ('PDF → pdf-parse  |  TXT/MD → buffer.toString()',  CYAN),
    ('chunker.js: 400-word / 60-word-overlap segments', PURPLE),
    ('embed.js: BM25 tokenize each chunk → JSON array', GREEN),
    ('INSERT doc metadata → kb_documents table',        BLUE),
    ('INSERT all chunks → kb_chunks table (batch)',     CYAN),
    ('Return { docId, chunkCount, name } to admin UI',  GREEN),
]
for i, (txt, acc) in enumerate(doc33):
    flow_node(sl, txt, 0.48, 1.52+i*0.70, 5.9, 0.50, acc, num=i+1)
    if i < 7: flow_arrow_v(sl, 3.23, 2.02+i*0.70)
code_box(sl, [
    'CREATE TABLE kb_documents (',
    '  id          TEXT PRIMARY KEY,',
    '  name        TEXT NOT NULL,',
    '  chunk_count INTEGER NOT NULL,',
    '  added_at    INTEGER NOT NULL',
    ');',
    '',
    'CREATE TABLE kb_chunks (',
    '  rowid     INTEGER PRIMARY KEY AUTOINCREMENT,',
    '  doc_id    TEXT NOT NULL',
    '            REFERENCES kb_documents(id)',
    '            ON DELETE CASCADE,',
    '  doc_name  TEXT NOT NULL,',
    '  text      TEXT NOT NULL,',
    '  embedding TEXT NOT NULL',
    '  -- JSON: ["token1","token2",...]',
    ');'],
    6.72, 1.52, 6.13, 5.0, lang='schema.sql')
snum(sl, 33)

# ── SLIDE 34  ERROR HANDLING ───────────────────────────────────────
sl = new_slide(prs, RGBColor(0x10, 0x06, 0x0b))
slide_header(sl, 'Error Handling', 'Error Handling & Fallback Flows', RED)
errs34 = [
    ('🔄','JSON Parse Fallback','Gemma returns JSON wrapped in markdown fences → strip before parse. On continued failure → Gemma 1.5 Flash in forced JSON mode.',RED),
    ('📭','Empty KB Results','KB-only mode + BM25 finds no relevant chunks → graceful "I don\'t have information on that in my knowledge base" message returned.',YELLOW),
    ('🌐','Gemma Web Search','Gemma models do not support the Google Search tool. webSearchWeight > 0 with Gemma → web search silently skipped. Admin warned in settings UI.',BLUE),
    ('🔊','TTS / STT Degradation','ElevenLabs API fails → UI shows error, text response still displayed. Browser lacks Web Speech → mic button hidden gracefully.',PURPLE),
]
for i, (icon, title, desc, acc) in enumerate(errs34):
    c = i % 2; rr = i // 2
    card(sl, 0.48+c*6.48, 1.52+rr*2.65, 6.15, 2.48, icon, title, desc, acc)
R(sl, 0.48, 6.6, 12.37, 0.62, fill=_dim(RED,10), line=RED, lw=0.5, radius=True)
T(sl, 'All API endpoints return consistent  {"error": "message"}  with appropriate HTTP codes (400 / 401 / 403 / 500). No stack traces exposed.',
  0.68, 6.68, 12.0, 0.46, size=11, color=LIGHT, align=PP_ALIGN.CENTER)
snum(sl, 34)

# ── SLIDE 35  COMPLETE SYSTEM FLOW ────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Full System', 'Complete System Interaction Overview', YELLOW)
tiers35 = [
    ('Client Tier',   BLUE,   '📱 React 18 SPA\n🔐 AuthContext (JWT)\n💬 ChatPage\n👤 ProfilePage\n🛡 AdminLayout\n🎤 Web Speech API\n🌗 Dark/Light Theme'),
    ('Server Tier',   PURPLE, '⚙️ Express REST API\n🔒 JWT Middleware\n📋 25+ Route Handlers\n📊 SQLite WAL\n🔍 BM25 RAG Engine\n✂️ Text Chunker\n🏗 Multi-Instance Spawn'),
    ('External APIs', GREEN,  '🤖 Google Gemma\n   (6 model variants)\n🔊 ElevenLabs TTS\n   (11 voices, 4 models)\n🏃 Google Fit API\n   (OAuth2, read-only)\n📄 pdf-parse library'),
]
for i, (title, acc, items) in enumerate(tiers35):
    R(sl, 0.42+i*4.33, 1.52, 4.06, 4.88, fill=_dim(acc,9), line=acc, lw=0.6, radius=True)
    R(sl, 0.42+i*4.33, 1.52, 0.08, 4.88, fill=acc)
    T(sl, title, 0.56+i*4.33, 1.62, 3.86, 0.34, size=12.5, bold=True, color=acc, align=PP_ALIGN.CENTER)
    hline(sl, 0.56+i*4.33, 2.04, 3.7, _dim(acc,4))
    T(sl, items, 0.56+i*4.33, 2.16, 3.86, 3.95, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 2: T(sl, '⟷', 4.36+i*4.33, 3.68, 0.48, 0.52, size=20, color=DIM, align=PP_ALIGN.CENTER)
sub35 = [('💾 Storage','SQLite WAL\n~12 MB typical',YELLOW),
         ('🔐 Auth','JWT + bcrypt\n7-day expiry',CYAN),
         ('☁️ Hosting','Render + Netlify\n1 GB Disk',RED),
         ('📈 Scaling','Multi-instance\nIsolated DBs',PURPLE)]
for i, (title, desc, acc) in enumerate(sub35):
    R(sl, 0.42+i*3.24, 6.52, 3.06, 0.72, fill=_dim(acc,10), line=_dim(acc,4), lw=0.5, radius=True)
    T(sl, title, 0.5+i*3.24, 6.58, 2.9, 0.28, size=10, bold=True, color=acc, align=PP_ALIGN.CENTER)
    T(sl, desc,  0.5+i*3.24, 6.86, 2.9, 0.32, size=9,  color=MUTED, align=PP_ALIGN.CENTER)
snum(sl, 35)

# ── SLIDE 36  SECTION 5 DIVIDER ────────────────────────────────────
section_slide(prs, 5, 'System Design',
    'Database schema  ·  REST API  ·  JWT lifecycle  ·  SQLite WAL  ·  Deployment',
    CYAN, 36)

# ── SLIDE 37  DB SCHEMA PART 1 ────────────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'DB Schema', 'Database Schema — Core Tables', CYAN)
code_box(sl, [
    'CREATE TABLE users (',
    '  id            TEXT PRIMARY KEY,  -- UUID',
    '  email         TEXT UNIQUE NOT NULL,',
    '  name          TEXT NOT NULL,',
    "  role          TEXT DEFAULT 'user',",
    '  password_hash TEXT NOT NULL,',
    '  login_count   INTEGER DEFAULT 0,',
    '  last_login    INTEGER,',
    "  note          TEXT DEFAULT '',",
    "  profile       TEXT DEFAULT '{}',  -- JSON",
    "  memories      TEXT DEFAULT '[]',  -- JSON",
    '  created_at    INTEGER NOT NULL',
    ');'],
    0.48, 1.52, 6.0, 4.12, lang='schema.sql')
code_box(sl, [
    'CREATE TABLE conversations (',
    '  id         TEXT PRIMARY KEY,',
    '  user_id    TEXT NOT NULL',
    '             REFERENCES users(id)',
    '             ON DELETE CASCADE,',
    "  title      TEXT DEFAULT 'New Chat',",
    '  created_at INTEGER NOT NULL,',
    '  updated_at INTEGER NOT NULL',
    ');',
    '',
    'CREATE TABLE messages (',
    '  id              TEXT PRIMARY KEY,',
    '  conversation_id TEXT NOT NULL',
    '    REFERENCES conversations(id)',
    '    ON DELETE CASCADE,',
    '  role      TEXT NOT NULL,',
    '  text      TEXT NOT NULL,',
    '  timestamp INTEGER NOT NULL,',
    '  refs      TEXT   -- JSON KB source refs',
    ');'],
    6.65, 1.52, 6.2, 4.12, lang='schema.sql')
info_box(sl, 0.48, 5.78, 12.37, 0.72,
         '📌  Design Notes',
         'All timestamps: Unix milliseconds (INTEGER).  UUIDs: crypto.randomUUID().  CASCADE deletes maintain referential integrity automatically.',
         CYAN)
snum(sl, 37)

# ── SLIDE 38  DB SCHEMA PART 2 ────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'DB Schema', 'Database Schema — Knowledge Base & Settings', GREEN)
code_box(sl, [
    'CREATE TABLE kb_documents (',
    '  id          TEXT PRIMARY KEY,',
    '  name        TEXT NOT NULL,',
    '  chunk_count INTEGER NOT NULL,',
    '  added_at    INTEGER NOT NULL',
    ');',
    '',
    'CREATE TABLE kb_chunks (',
    '  rowid     INTEGER PRIMARY KEY AUTOINCREMENT,',
    '  doc_id    TEXT NOT NULL',
    '            REFERENCES kb_documents(id)',
    '            ON DELETE CASCADE,',
    '  doc_name  TEXT NOT NULL,',
    '  text      TEXT NOT NULL,',
    '  embedding TEXT NOT NULL',
    '  -- JSON array of BM25 tokens',
    "  -- e.g. ['diabetes','management']",
    ');'],
    0.48, 1.52, 6.0, 4.12, lang='schema.sql')
code_box(sl, [
    'CREATE TABLE settings (',
    '  key   TEXT PRIMARY KEY,',
    '  value TEXT NOT NULL   -- JSON blob',
    ');',
    '',
    '-- Single row "botSettings" (full example):',
    '{',
    '  "model":              "gemma-3-flash-preview",',
    '  "systemPrompt":       "You are CIRA...",',
    '  "style":              "balanced",',
    '  "ragTopK":            5,',
    '  "ragWeight":          80,',
    '  "ownKnowledgeWeight": 60,',
    '  "webSearchWeight":    40,',
    '  "defaultTheme":       "dark",',
    '  "ttsEnabled":         true,',
    '  "defaultVoiceId":     "21m00Tcm....",',
    '  "defaultTtsModel":    "eleven_turbo_v2_5",',
    '  "ttsStability":       0.5,',
    '  "ttsSimilarityBoost": 0.75',
    '}'],
    6.65, 1.52, 6.2, 4.12, lang='schema.sql')
info_box(sl, 0.48, 5.78, 12.37, 0.72,
         '⚡  Live Settings',
         'Admin changes take effect on the very next chat request — no server restart or redeploy required.',
         GREEN)
snum(sl, 38)

# ── SLIDE 39  REST API ─────────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'REST API', 'REST API Architecture — 25+ Endpoints', BLUE)
table(sl, ['Endpoint','Method','Auth','Purpose'],
      [['/api/auth/signup','POST','Public','Register user; returns JWT'],
       ['/api/auth/login','POST','Public','Authenticate; returns JWT'],
       ['/api/auth/me','GET','User JWT','Verify token; return user object'],
       ['/api/conversations','GET / POST','User JWT','List or create conversations'],
       ['/api/conversations/:id','PATCH / DELETE','User JWT','Rename or delete conversation'],
       ['/api/conversations/:id/messages','GET','User JWT','Fetch message history for thread'],
       ['/api/chat','POST','User JWT','Send message → AI response (core)'],
       ['/api/title','POST','User JWT','Auto-generate conversation title via Gemma'],
       ['/api/health','GET','Public','Health check for Render.com monitoring']],
      0.48, 1.52, 5.9, 4.35, BLUE, col_widths=[2.6, 1.15, 1.15, 0.0])
table(sl, ['Endpoint','Method','Auth','Purpose'],
      [['/api/user/profile','GET / PUT','User JWT','Read or update health profile fields'],
       ['/api/user/memories','GET / DELETE','User JWT','List memories or delete by ID'],
       ['/api/user/memories/extract','POST','User JWT','Trigger Gemma memory extraction'],
       ['/api/fitness/auth-url','GET','User JWT','Get Google Fit OAuth2 redirect URL'],
       ['/api/fitness/callback','GET','Public','Handle OAuth2 code exchange'],
       ['/api/fitness/refresh','POST','User JWT','Fetch fresh 7-day fitness data'],
       ['/api/fitness/disconnect','DELETE','User JWT','Revoke Fit connection + clear tokens'],
       ['/api/tts','POST','User JWT','Convert text to MP3 via ElevenLabs'],
       ['/api/admin/users','GET','Admin JWT','List all users with stats'],
       ['/api/admin/settings','GET / PUT','Admin JWT','Read or update bot configuration'],
       ['/api/documents','GET/POST/DELETE','Admin JWT','Manage knowledge base documents']],
      6.65, 1.52, 6.2, 5.12, PURPLE, col_widths=[2.4, 1.15, 1.15, 0.0])
info_box(sl, 0.48, 5.76, 12.37, 0.82,
         'REST Design Principles',
         'Resource-oriented URLs  ·  HTTP verbs for actions  ·  Consistent {"error":"msg"} error shape  ·  Stateless JWT auth  ·  JSON + MP3 content types',
         BLUE)
snum(sl, 39)

# ── SLIDE 40  JWT DESIGN ──────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'JWT Design', 'JWT Authentication Design', YELLOW)
code_box(sl, [
    '// Header (base64url-encoded)',
    '{ "alg": "HS256", "typ": "JWT" }',
    '',
    '// Payload (base64url-encoded)',
    '{',
    '  "uid":   "user-uuid-here",',
    '  "email": "user@example.com",',
    '  "name":  "John Doe",',
    '  "role":  "user",',
    '  "iat":   1716890400,',
    '  "exp":   1717495200   // +7 days',
    '}',
    '',
    '// Signature (HMAC-SHA256)',
    'HMAC(',
    '  base64(header) + "." + base64(payload),',
    '  JWT_SECRET',
    ')'],
    0.48, 1.52, 6.0, 5.6, lang='JWT')
arch_box(sl, 6.72, 1.52, 6.13, '🔐', 'Stateless Authentication', 'No session store needed. Server only needs JWT_SECRET to verify. Scales horizontally without shared session state.', YELLOW)
arch_box(sl, 6.72, 2.48, 6.13, '👑', 'Role-Based Access Control', 'role claim in payload gates admin routes. requireAdmin checks role==="admin" after requireAuth JWT verification.', PURPLE)
arch_box(sl, 6.72, 3.44, 6.13, '💾', 'Client-Side Storage', 'Token stored in localStorage as "chatbot_token". AuthContext reads on app load. 7-day expiry for session continuity.', BLUE)
arch_box(sl, 6.72, 4.40, 6.13, '⏱️', 'Token Lifecycle', 'Issued on login/signup → sent as Bearer token on every request → 7-day expiry → user must re-authenticate.', CYAN)
snum(sl, 40)

# ── SLIDE 41  SQLITE WAL ──────────────────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'SQLite WAL', 'SQLite WAL Mode — Why and How', CYAN)
bullets(sl, [
    ('WAL = Write-Ahead Logging',     'Writes go to a separate .db-wal log file first, then checkpointed to the main .db file asynchronously.'),
    ('Readers never block writers',   'Multiple SELECT queries can run concurrently with INSERT/UPDATE — no shared lock on the main file during writes.'),
    ('Concurrent HTTP requests',      'Multiple users chatting simultaneously share the DB file safely. Handles ~50 concurrent users per instance without degradation.'),
    ('Faster write throughput',       'Chat message inserts (2 rows per chat turn) complete in 1–2 ms. WAL avoids the exclusive lock of rollback-journal mode.'),
    ('3 files created automatically', 'chatbot.db (main data) · chatbot.db-wal (pending writes) · chatbot.db-shm (shared memory index). All managed automatically.'),
    ('Automatic crash recovery',      'Uncommitted WAL entries are rolled back automatically on the next open() after a crash — zero data corruption risk.'),
    ('foreign_keys = ON',             'SQLite foreign key enforcement is opt-in and enabled at connection time. CASCADE deletes maintain referential integrity.'),
    ('Zero operational overhead',     'Embedded in the Node.js process. No server to start, no port to configure, no version to manage.'),
], 0.48, 1.52, 5.9, 5.38, CYAN)
code_box(sl, [
    '// db.js — initialization',
    '',
    "const db = new Database(dbPath);",
    "db.pragma('journal_mode = WAL');",
    "db.pragma('foreign_keys = ON');",
    '',
    '// WAL creates 3 files:',
    '//   chatbot.db         (main)',
    '//   chatbot.db-wal     (write-ahead log)',
    '//   chatbot.db-shm     (shared memory)'],
    0.48, 5.38, 5.9, 1.8, lang='db.js')
table(sl, ['','SQLite WAL','PostgreSQL'],
      [['Setup','Zero config','Server required'],
       ['Concurrent R/W','✓  (WAL mode)','✓  (MVCC)'],
       ['File size','~12 MB','Larger overhead'],
       ['Monthly cost','$0','$$$'],
       ['Multi-host scale','Single node','Multi-host'],
       ['Right for CIRA?','✅  Perfect','Overkill']],
      6.72, 1.52, 6.13, 4.65, CYAN)
snum(sl, 41)

# ── SLIDE 42  FRONTEND ARCHITECTURE ───────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Frontend Design', 'Frontend Application Architecture', BLUE)
T(sl, 'Page Structure', 0.48, 1.52, 5.9, 0.30, size=11.5, bold=True, color=BLUE)
arch_box(sl, 0.48, 1.88, 5.9, '🔑', '/ — LoginPage.jsx', 'JWT-authenticated entry point. Redirects to /chat if logged in. Sign-up and login forms with client-side validation.', BLUE, h=0.88)
arch_box(sl, 0.48, 2.86, 5.9, '💬', '/chat — ChatPage.jsx', 'Main interface: sidebar with conversations, message thread, input bar with STT button. Markdown rendering via react-markdown.', PURPLE, h=0.88)
arch_box(sl, 0.48, 3.84, 5.9, '👤', '/profile — ProfilePage.jsx', 'Health profile form, AI memory manager, Google Fit connect button, 7-day fitness snapshot display.', GREEN, h=0.88)
arch_box(sl, 0.48, 4.82, 5.9, '🛡️', '/admin/* — AdminLayout.jsx', 'Dashboard, Users, Chat History, KB Manager, Settings — tabbed admin interface with admin-only access guard.', YELLOW, h=0.88)
T(sl, 'State Management — Context API', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=PURPLE)
code_box(sl, [
    '// AuthContext.jsx — single global context',
    'const AuthContext = createContext();',
    '',
    'export function AuthProvider({ children }) {',
    '  const [user, setUser] = useState(null);',
    '  const [token, setToken] = useState(',
    "    localStorage.getItem('chatbot_token')",
    '  );',
    '',
    '  // Auto-verify JWT on mount',
    '  useEffect(() => {',
    '    if (token) verifyToken(token)',
    '      .then(setUser);',
    '  }, []);',
    '',
    '  return (',
    '    <AuthContext.Provider',
    '      value={{ user, token, login, logout }}',
    '    >',
    '      {children}',
    '    </AuthContext.Provider>',
    '  );',
    '}'],
    6.72, 1.88, 6.13, 5.35, lang='AuthContext.jsx')
snum(sl, 42)

# ── SLIDE 43  BACKEND ARCHITECTURE ────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Backend Design', 'Backend Architecture', PURPLE)
table(sl, ['File / Module','Responsibility & Key Details'],
      [['server.js','Main Express app — all 25+ route handlers, Gemma API integration, TTS, Fit. ~966 lines.'],
       ['db.js','SQLite init: WAL mode, foreign_keys=ON, schema creation, all prepared statements.'],
       ['auth.js','requireAuth middleware (jwt.verify) + requireAdmin role check for admin routes.'],
       ['rag/embed.js','BM25 tokenizer: lowercase, 47-stopword filter, token-frequency map → JSON array.'],
       ['rag/chunker.js','400-word / 60-word-overlap text splitter. Returns [{text, index}] array.'],
       ['rag/store.js','BM25 search engine: IDF precompute, per-chunk scoring, Top-K sort, returns with source names.'],
       ['start-instances.js','Reads instances.json, spawns isolated Node.js child processes with per-instance env vars.'],
       ['instances.json','Config array: [{name, port, dataDir}] — one entry per chatbot instance deployed.']],
      0.48, 1.52, 5.9, 3.8, PURPLE, col_widths=[2.15, 3.75])
T(sl, 'Request Middleware Chain', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=BLUE)
chain43 = [
    ('CORS  (allow FRONTEND_URL only)',                BLUE),
    ('express.json()  body parser',                   CYAN),
    ('requireAuth  — JWT verify',                     YELLOW),
    ('requireAdmin  — role check  (admin routes only)',PURPLE),
    ('Route Handler → SQLite / Gemma / ElevenLabs',  GREEN),
    ('res.json()  response to client',                BLUE),
]
for i, (txt, acc) in enumerate(chain43):
    flow_node(sl, txt, 6.72, 1.88+i*0.86, 6.13, 0.58, acc, num=i+1)
    if i < 5: flow_arrow_v(sl, 9.59, 2.46+i*0.86)
snum(sl, 43)

# ── SLIDE 44  DEPLOYMENT SYSTEM DESIGN ───────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Deployment', 'Deployment System Design', CYAN)
code_box(sl, [
    '# render.yaml',
    'services:',
    '  - type: web',
    '    name: cira-backend',
    '    runtime: node',
    '    buildCommand: "cd backend && npm i"',
    '    startCommand: "node server.js"',
    '    disk:',
    '      name: data',
    '      mountPath: /app/data',
    '      sizeGB: 1',
    '    healthCheckPath: /api/health'],
    0.48, 1.52, 6.0, 3.3, lang='render.yaml')
code_box(sl, [
    '# netlify.toml',
    '[build]',
    '  base    = "frontend"',
    '  command = "npm run build"',
    '  publish = "build"',
    '',
    '[[redirects]]',
    '  from   = "/*"',
    '  to     = "/index.html"',
    '  status = 200   # SPA routing'],
    0.48, 4.96, 6.0, 2.38, lang='netlify.toml')
info_box(sl, 6.72, 1.52, 6.13, 1.0, '☁️  Netlify CDN (Frontend)',
         'React SPA served globally from edge nodes. Instant cache invalidation on push. Zero cold-start.', BLUE)
info_box(sl, 6.72, 2.62, 6.13, 1.0, '⚙️  Render.com (Backend)',
         'Node.js + 1 GB persistent disk for SQLite. Auto-deploy on GitHub push. Zero-downtime deploys.', PURPLE)
info_box(sl, 6.72, 3.72, 6.13, 1.0, '🤖  External APIs',
         'Google Gemma · ElevenLabs TTS · Google Fit — all HTTPS + API key auth. No vendor lock-in for AI model.', GREEN)
info_box(sl, 6.72, 4.82, 6.13, 1.0, '📈  Scaling Path',
         'SQLite WAL for small-medium. Clear migration to PostgreSQL + multiple instances for enterprise scale.', YELLOW)
snum(sl, 44)

# ── SLIDE 45  SECTION 6 DIVIDER ────────────────────────────────────
section_slide(prs, 6, 'Architecture & Advantages',
    'Design patterns  ·  Cost efficiency  ·  Security  ·  Scalability  ·  Why CIRA wins',
    RED, 45)

# ── SLIDE 46  HIGH-LEVEL ARCHITECTURE DIAGRAM ─────────────────────
sl = new_slide(prs)
slide_header(sl, 'Architecture', 'High-Level Architecture Diagram', RED)
tiers46 = [
    ('Client Tier',   BLUE,   '📱 React 18 SPA\n🔐 AuthContext (JWT)\n💬 ChatPage\n👤 ProfilePage\n🛡 AdminLayout\n🎤 Web Speech API\n🌗 Dark/Light Theme'),
    ('Server Tier',   PURPLE, '⚙️ Express REST API\n🔒 JWT Middleware\n📋 25+ Routes\n📊 SQLite WAL\n🔍 BM25 RAG Engine\n✂️ Text Chunker\n🏗 Multi-Instance Spawn'),
    ('External APIs', GREEN,  '🤖 Google Gemma\n   (6 model variants)\n🔊 ElevenLabs TTS\n   (11 voices, 4 models)\n🏃 Google Fit API\n   (OAuth2, read-only)\n📄 pdf-parse library'),
]
for i, (title, acc, items) in enumerate(tiers46):
    R(sl, 0.42+i*4.33, 1.52, 4.06, 4.95, fill=_dim(acc,9), line=acc, lw=0.6, radius=True)
    R(sl, 0.42+i*4.33, 1.52, 0.08, 4.95, fill=acc)
    T(sl, title, 0.60+i*4.33, 1.62, 3.82, 0.34, size=12.5, bold=True, color=acc, align=PP_ALIGN.CENTER)
    hline(sl, 0.60+i*4.33, 2.04, 3.72, _dim(acc,4))
    T(sl, items, 0.60+i*4.33, 2.16, 3.82, 4.0, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 2: T(sl, '⟷', 4.35+i*4.33, 3.7, 0.50, 0.55, size=20, color=DIM, align=PP_ALIGN.CENTER)
T(sl, 'Client  ←→  HTTPS REST API  ←→  SQLite     |     Server  ←→  Gemma / ElevenLabs / Google Fit',
  0.48, 6.6, 12.37, 0.32, size=9.5, color=DIM, align=PP_ALIGN.CENTER)
snum(sl, 46)

# ── SLIDE 47  DESIGN PATTERNS ─────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Design Patterns', 'Architectural Design Patterns', RED)
patterns47 = [
    ('🏛','Client-Server (3-Tier)','React (presentation) ↔ Express (business logic) ↔ SQLite (data). Clean separation enables independent deployment and scaling.',BLUE),
    ('🔗','REST Stateless API','All interactions via HTTP verbs on resource URIs. Stateless requests enabled by JWT, eliminating server-side session state.',PURPLE),
    ('🧩','RAG Pattern','Domain knowledge retrieved at query time and injected into LLM context. Eliminates hallucination risk for domain-specific questions.',CYAN),
    ('🔐','Middleware Chain (AOP)','Cross-cutting concerns (auth, CORS, body parsing) separated into Express middleware. Business logic stays clean.',GREEN),
    ('🏭','Multi-Tenant (Process Isolation)','Separate OS processes per instance, each with isolated SQLite file. True data isolation without containers.',YELLOW),
    ('🎯','Strategy Pattern (Knowledge)','RAG, own knowledge, and web search are independently togglable strategies with configurable weight percentages.',RED),
]
for i, (icon, title, desc, acc) in enumerate(patterns47):
    c = i % 2; rr = i // 2
    R(sl, 0.48+c*6.48, 1.52+rr*1.86, 6.15, 1.72, fill=_dim(acc,10), line=acc, lw=0.5, radius=True)
    R(sl, 0.48+c*6.48, 1.52+rr*1.86, 0.06, 1.72, fill=acc)
    T(sl, icon+' '+title, 0.68+c*6.48, 1.64+rr*1.86, 5.75, 0.32, size=12.5, bold=True, color=acc)
    T(sl, desc, 0.68+c*6.48, 2.00+rr*1.86, 5.75, 1.08, size=10.5, color=MUTED)
snum(sl, 47)

# ── SLIDE 48  ADVANTAGE: NO VECTOR EMBEDDINGS ─────────────────────
sl = new_slide(prs, RGBColor(0x05, 0x12, 0x0c))
slide_header(sl, 'Advantage #1', 'Zero Vector Infrastructure', GREEN)
for i, (num, lbl, col) in enumerate([('$0','Embedding API Cost',GREEN),
                                      ('<5ms','BM25 Search Latency',CYAN),
                                      ('12 MB','Typical DB Size',YELLOW)]):
    stat_box(sl, 0.48+i*4.3, 1.52, 4.05, 1.62, num, lbl, col)
R(sl, 0.48, 3.3, 5.9, 3.58, fill=_dim(RED,10), line=RED, lw=0.6, radius=True)
R(sl, 0.48, 3.3, 0.06, 3.58, fill=RED)
T(sl, '❌  Traditional RAG Stack', 0.68, 3.40, 5.6, 0.30, size=11.5, bold=True, color=RED)
bullets(sl, [
    'OpenAI Embedding API: ~$0.13 / million tokens',
    'Pinecone / Weaviate: $70–250 / month',
    'Requires separate async embedding pipeline',
    'Slow cold-start and external API dependency',
], 0.68, 3.78, 5.6, 2.8, RED, size=11)
R(sl, 6.72, 3.3, 6.13, 3.58, fill=_dim(GREEN,10), line=GREEN, lw=0.6, radius=True)
R(sl, 6.72, 3.3, 0.06, 3.58, fill=GREEN)
T(sl, '✅  CIRA BM25 Advantage', 6.92, 3.40, 5.9, 0.30, size=11.5, bold=True, color=GREEN)
bullets(sl, [
    'Pure JavaScript — zero external API calls',
    'SQLite storage: free, embedded, fully portable',
    'Instant synchronous indexing (< 5 ms)',
    'Fully explainable results (keyword matching)',
], 6.92, 3.78, 5.9, 2.8, GREEN, size=11)
snum(sl, 48)

# ── SLIDE 49  ADVANTAGE: MODULAR KNOWLEDGE ────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Advantage #2', 'Modular, Weighted Knowledge Sources', CYAN)
T(sl, 'Admins independently tune 3 knowledge source weights from 0–100%, composing custom strategies per deployment:',
  0.48, 1.52, 12.37, 0.42, size=12, color=MUTED)
srcs49 = [
    ('📚','KB Documents  (RAG)',    'Upload domain docs (PDF/TXT/MD). Ground answers strictly in your content. Set to 100% for compliance-sensitive deployments.',BLUE),
    ('🧠',"Gemma's Own Knowledge", "Leverage Gemma's training data for questions outside the KB. Set to 0% for strict KB-only mode.",PURPLE),
    ('🌐','Live Web Search',        'Real-time Google Search via Gemma tool use. Always-current answers. Disable for air-gapped environments.',CYAN),
]
for i, (icon, title, desc, acc) in enumerate(srcs49):
    R(sl, 0.48+i*4.3, 2.1, 4.05, 3.52, fill=_dim(acc,9), line=acc, lw=0.6, radius=True)
    R(sl, 0.48+i*4.3, 2.1, 0.06, 3.52, fill=acc)
    T(sl, icon, 0.65+i*4.3, 2.22, 3.82, 0.55, size=24, align=PP_ALIGN.CENTER)
    T(sl, title, 0.65+i*4.3, 2.85, 3.82, 0.35, size=12, bold=True, color=acc, align=PP_ALIGN.CENTER)
    hline(sl, 0.7+i*4.3, 3.28, 3.65, _dim(acc,4))
    T(sl, desc, 0.65+i*4.3, 3.42, 3.82, 1.9, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
R(sl, 0.48, 5.76, 12.37, 0.78, fill=_dim(BLUE,10), line=STROKE, lw=0.5, radius=True)
T(sl, 'Presets:  Strict (KB=100, Own=0, Web=0)   ·   Balanced (KB=80, Own=60, Web=40)   ·   General (KB=0, Own=100, Web=100)',
  0.68, 5.95, 12.0, 0.42, size=11, color=MUTED, align=PP_ALIGN.CENTER)
snum(sl, 49)

# ── SLIDE 50  ADVANTAGE: SMART PERSONALIZATION ────────────────────
sl = new_slide(prs)
slide_header(sl, 'Advantage #3', 'Smart Personalization Without Manual Prompting', GREEN)
T(sl, "CIRA's personalization stack works entirely automatically — users never need to re-explain themselves:",
  0.48, 1.52, 5.9, 0.48, size=12, color=MUTED)
for i, (title, desc, acc) in enumerate([
    ('User fills profile once',   'Health data stored permanently and auto-injected into every conversation context.',GREEN),
    ('AI extracts facts automatically','Memories extracted from chat history via Gemma — system learns preferences without forms.',CYAN),
    ('Fitness data syncs on demand','One OAuth2 connection to Google Fit. Steps, HR, weight always available in context.',BLUE),
    ('All assembled per-request', 'Profile + memories + fitness + KB + history = rich systemInstruction in < 20 ms.',YELLOW),
]):
    R(sl, 0.48, 2.12+i*1.28, 0.62, 0.62, fill=_dim(acc,4), line=acc, lw=0.6, radius=True)
    T(sl, str(i+1), 0.48, 2.12+i*1.28, 0.62, 0.62, size=15, bold=True, color=acc, align=PP_ALIGN.CENTER)
    R(sl, 1.2, 2.12+i*1.28, 5.18, 0.62, fill=_dim(acc,10), line=_dim(acc,4), lw=0.4, radius=True)
    T(sl, title, 1.34, 2.17+i*1.28, 4.9, 0.28, size=11.5, bold=True, color=WHITE)
    T(sl, desc,  1.34, 2.48+i*1.28, 4.9, 0.62, size=9.5, color=MUTED)
table(sl, ['Feature','CIRA','Generic Chatbots'],
      [['User profile','✅  Stored & injected','❌'],
       ['Learned memories','✅  Auto-extracted','❌'],
       ['Fitness data','✅  Live via OAuth2','❌'],
       ['Conversation history','✅  Persistent','Session only'],
       ['Tone / style control','✅  Admin config','❌']],
      6.72, 1.52, 6.13, 4.08, GREEN, col_widths=[2.5, 1.8, 1.83])
snum(sl, 50)

# ── SLIDE 51  ADVANTAGE: COST EFFICIENCY ──────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Advantage #4', 'Cost Efficiency — Run for Under $20/Month', YELLOW)
table(sl, ['Component','CIRA Cost','Typical Alternative'],
      [['Hosting (Render)','$7 / month','$20–50'],
       ['CDN (Netlify)','$0  (free tier)','$15–20'],
       ['Vector Database','$0  (SQLite)','$70–250'],
       ['Embedding API','$0  (BM25)','$5–50'],
       ['AI API (Gemma)','Pay-per-use','Pay-per-use'],
       ['TTS (ElevenLabs)','$5 starter','$5+'],
       ['Total fixed infra','~$12 / month','$110–370']],
      0.48, 1.52, 5.9, 4.08, YELLOW, col_widths=[2.4, 1.5, 2.0])
T(sl, 'Cost-Saving Architecture Decisions', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=CYAN)
bullets(sl, [
    ('BM25 over vectors',      'Saves $70–250/month in vector DB costs alone'),
    ('SQLite over PostgreSQL', 'Zero database hosting cost and zero maintenance'),
    ('Context API',            'No state management library overhead'),
    ('Web Speech API',         'Zero STT API cost — native browser capability'),
    ('Gemma Flash',           '10–50× cheaper per token than GPT-4 equivalent'),
    ('Multi-instance on one server', 'One Render instance serves N chatbots in parallel'),
], 6.72, 1.88, 6.13, 3.72, YELLOW)
snum(sl, 51)

# ── SLIDE 52  ADVANTAGE: SCALABILITY ─────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Advantage #5', 'Scalability & Reliability', BLUE)
scalability52 = [
    ('⚡','Stateless REST API','JWT auth eliminates server-side session state. Any number of Node.js processes can serve the same JWT — horizontal scaling with no shared session store.',BLUE),
    ('🔀','Multi-Instance Process Isolation','Each chatbot instance runs as an isolated OS process with its own SQLite file. Instance A crashing does not affect Instance B.',PURPLE),
    ('📦','SQLite WAL Concurrent I/O','WAL mode handles concurrent HTTP requests without blocking. Multiple users chatting simultaneously share the DB safely — ~50 concurrent users per instance.',CYAN),
    ('🌐','CDN-Served Frontend','React SPA is a static build served from Netlify\'s global edge network. Frontend scales infinitely regardless of backend load.',GREEN),
]
for i, (icon, title, desc, acc) in enumerate(scalability52):
    c = i % 2; rr = i // 2
    R(sl, 0.48+c*6.48, 1.52+rr*2.72, 6.15, 2.56, fill=_dim(acc,10), line=acc, lw=0.5, radius=True)
    R(sl, 0.48+c*6.48, 1.52+rr*2.72, 0.06, 2.56, fill=acc)
    T(sl, icon+' '+title, 0.68+c*6.48, 1.64+rr*2.72, 5.76, 0.32, size=12.5, bold=True, color=acc)
    T(sl, desc, 0.68+c*6.48, 2.02+rr*2.72, 5.76, 1.85, size=10.5, color=MUTED)
snum(sl, 52)

# ── SLIDE 53  ADVANTAGE: SECURITY ────────────────────────────────
sl = new_slide(prs, RGBColor(0x10, 0x06, 0x0b))
slide_header(sl, 'Advantage #6', 'Security Architecture', RED)
sec53 = [
    ('🔑','Password Security','bcryptjs with 10 salt rounds. Never stored in plaintext. Computationally expensive to brute-force.',RED),
    ('🎫','JWT Signed Tokens','HMAC-SHA256 signed. Tamper-evident. 7-day expiry. Role claims verified server-side on every request.',YELLOW),
    ('🚧','CORS Protection','Only FRONTEND_URL allowed as CORS origin. Blocks unauthorized cross-origin API calls.',BLUE),
    ('🔒','OAuth2 Read-Only','Google Fit requests only read scopes. Tokens stored server-side. Never exposed to frontend.',PURPLE),
    ('🛡️','Role-Based Access','Admin routes double-protected: requireAuth (JWT valid) + requireAdmin (role="admin"). Layered defence.',GREEN),
    ('📝','Safe Markdown','react-markdown sanitizes HTML in AI responses by default. No XSS injection through AI-generated content.',CYAN),
]
for i, (icon, title, desc, acc) in enumerate(sec53):
    c = i % 3; rr = i // 3
    card(sl, 0.48+c*4.3, 1.52+rr*2.72, 4.05, 2.55, icon, title, desc, acc)
snum(sl, 53)

# ── SLIDE 54  ADVANTAGE: DEVELOPER EXPERIENCE ────────────────────
sl = new_slide(prs)
slide_header(sl, 'Advantage #7', 'Developer Experience & Extensibility', PURPLE)
bullets(sl, [
    ('Zero build config',       'CRA frontend, plain Node.js backend. No Webpack config, no TypeScript migration needed.'),
    ('No ORM needed',           'Raw better-sqlite3 with synchronous API — easy to reason about, debug, and extend.'),
    ('Model switching',         'Change Gemma model in admin settings UI — zero code change required.'),
    ('Behavior without redeploy','System prompt, temperature, KB weights, TTS defaults — all in database settings.'),
    ('Self-contained install',  'Single npm install in /frontend and /backend — no monorepo tooling.'),
    ('Instances in JSON',       'Add a new chatbot instance by adding one JSON object to instances.json.'),
    ('CSS design tokens',       'Theme customisation via CSS custom properties — no Tailwind config needed.'),
], 0.48, 1.52, 5.9, 5.68, PURPLE)
T(sl, 'Start in 4 Steps', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=BLUE)
code_box(sl, [
    '# 1. Clone and install',
    'git clone https://github.com/.../cira',
    'cd backend   && npm install',
    'cd ../frontend && npm install',
    '',
    '# 2. Configure environment',
    'cp backend/.env.example backend/.env',
    '# → Set GEMINI_API_KEY, JWT_SECRET',
    '',
    '# 3. Start all instances',
    'node start-instances.js',
    '',
    '# 4. Start frontend',
    'cd frontend && npm start',
    '',
    '# Ready at http://localhost:3000 🚀'],
    6.72, 1.88, 6.13, 5.32, lang='Terminal')
snum(sl, 54)

# ── SLIDE 55  ADVANTAGE: ADMIN-FIRST ─────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Advantage #8', 'Admin-First: Zero-Code Bot Configuration', YELLOW)
T(sl, 'Every aspect of CIRA behaviour can be changed by a non-technical admin without touching any code:',
  0.48, 1.52, 12.37, 0.42, size=12, color=MUTED)
configs55 = [
    ('🤖','Switch AI Model',    'Dropdown: Gemma Flash → Pro → Gemma. Takes effect on the next request.',BLUE),
    ('📝','Edit System Prompt', 'Change bot persona, topic scope, tone, language via a text area in admin settings.',PURPLE),
    ('🌡️','Response Style',    'Toggle: Precise (0.2) / Balanced (0.7) / Creative (1.2) temperature presets.',CYAN),
    ('📚','Update Knowledge Base','Upload new PDFs or delete outdated docs. Available instantly after upload.',GREEN),
    ('🔊','Configure TTS',      'Set default voice, model, stability, similarity boost for all users.',YELLOW),
    ('⚖️','Knowledge Weights',  'Sliders for RAG, own knowledge, and web search weights (0–100 each).',RED),
]
for i, (icon, title, desc, acc) in enumerate(configs55):
    c = i % 3; rr = i // 3
    card(sl, 0.48+c*4.3, 2.12+rr*2.45, 4.05, 2.28, icon, title, desc, acc)
snum(sl, 55)

# ── SLIDE 56  MULTI-MODEL SUPPORT ────────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'Model Flexibility', 'Multi-Model Support', CYAN)
T(sl, 'CIRA supports all major Gemma generations, enabling cost-performance trade-offs optimised per deployment:',
  0.48, 1.52, 12.37, 0.42, size=12, color=MUTED)
table(sl, ['Model ID','Speed','Quality','Cost / 1M tok','Web Search','Ideal Deployment'],
      [['gemma-3-flash-preview  ⭐','Fastest','Excellent','~$0.075','✅  Built-in','Default — best balance of speed and quality for all use cases'],
       ['gemma-2.5-pro',  'Moderate','Best in class','~$1.25','✅  Built-in','Deep reasoning, complex multi-step medical or legal Q&A'],
       ['gemma-2.5-flash','Fast',    'Very Good','~$0.15','✅  Built-in','High-volume production deployments needing speed + quality'],
       ['gemma-2.0-flash','Fast',    'Good',   '~$0.10','✅  Built-in','Budget-optimised or legacy compatibility use cases'],
       ['gemma-4 (open-weight)','Variable','Good','Free (self-host)','❌  Not supported','Air-gapped, privacy-first, or offline-capable deployments']],
      0.48, 2.05, 12.37, 3.75, CYAN,
      col_widths=[2.8,1.0,1.15,1.25,1.3,4.87])
info_box(sl, 0.48, 5.96, 12.37, 0.78,
         'Admin Model Switching & Fallback',
         'Admin switches model via a settings dropdown with zero code changes. Selecting Gemma 4 (open-weight) automatically hides the web-search slider in the UI and silently disables grounding at runtime. Memory extraction uses a hardcoded Gemma Flash fallback in JSON mode for parse-failure robustness.',
         CYAN)
snum(sl, 56)

# ── SLIDE 57  PROGRESSIVE ENHANCEMENT ───────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Design Principle', 'Progressive Enhancement — Graceful Degradation', PURPLE)
T(sl, 'Every optional feature degrades gracefully — the core chat experience is never broken by unavailable optional services:',
  0.48, 1.52, 12.37, 0.42, size=12, color=MUTED)
table(sl, ['Feature','If Unavailable — Fallback Behaviour'],
      [['ElevenLabs TTS',    'TTS button hidden; text response still displayed normally'],
       ['Web Speech (STT)',  'Mic button hidden; keyboard input works normally'],
       ['Google Fit',        'Fitness section skipped; profile data still injected into context'],
       ['KB Documents',      'RAG weight ignored; own knowledge and web search used instead'],
       ['Web Search',        'Silently disabled on Gemma; KB + own knowledge used'],
       ['User Profile',      'System prompt only; no personalization context injected']],
      0.48, 2.05, 12.37, 3.1, PURPLE, col_widths=[2.8, 9.57])
R(sl, 0.48, 5.28, 5.9, 1.08, fill=_dim(GREEN,10), line=GREEN, lw=0.5, radius=True)
R(sl, 0.48, 5.28, 0.06, 1.08, fill=GREEN)
T(sl, '✅  Core Experience Guarantee', 0.68, 5.38, 5.7, 0.30, size=11.5, bold=True, color=GREEN)
T(sl, 'As long as GEMINI_API_KEY and JWT_SECRET are set, the basic text chat experience always works.',
  0.68, 5.74, 5.7, 0.55, size=10, color=MUTED)
R(sl, 6.65, 5.28, 6.2, 1.08, fill=_dim(BLUE,10), line=BLUE, lw=0.5, radius=True)
R(sl, 6.65, 5.28, 0.06, 1.08, fill=BLUE)
T(sl, '🚀  Minimum Viable Deployment', 6.85, 5.38, 5.9, 0.30, size=11.5, bold=True, color=BLUE)
T(sl, 'Set GEMINI_API_KEY + JWT_SECRET → get a fully functional, authenticated, multi-conversation chatbot immediately.',
  6.85, 5.74, 5.9, 0.55, size=10, color=MUTED)
snum(sl, 57)

# ── SLIDE 58  COMPARISON WITH ALTERNATIVES ───────────────────────
sl = new_slide(prs)
slide_header(sl, 'Comparison', 'CIRA vs. Alternatives', BLUE)
table(sl, ['Feature','CIRA','ChatGPT API','LangChain','Botpress'],
      [['Self-hosted',         '✅  Full control','❌  Cloud only','✅','Partial'],
       ['Vector DB needed',    '❌  BM25 only',  'Yes','Yes','Yes'],
       ['User personalization','✅  Deep',        'Manual only','Custom build','Limited'],
       ['Google Fit',          '✅  Built-in',    '❌','Custom build','❌'],
       ['Admin dashboard',     '✅  Full',        '❌','❌','✅'],
       ['AI memory extraction','✅  Auto (Gemma)','Manual API','Custom build','❌'],
       ['Monthly cost (infra)','~$12',            '$100+','$50–200','$200+'],
       ['Setup complexity',    'Low (2 installs)','Medium','High','High']],
      0.48, 1.52, 12.37, 5.15, BLUE,
      col_widths=[2.8, 2.3, 2.0, 2.0, 3.27])
snum(sl, 58)

# ── SLIDE 59  PERFORMANCE ─────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Performance', 'Performance Characteristics', GREEN)
table(sl, ['Operation','P50 Latency','P99 Latency','Notes'],
      [['JWT verify (requireAuth)',              '< 0.5 ms','< 1 ms',  'Pure CPU — HMAC-SHA256 hash. No DB query.'],
       ['SQLite read: user profile + memories', '3–8 ms',  '12 ms',   'WAL mode — single indexed SELECT by UUID.'],
       ['SQLite read: last 10 messages',        '5–10 ms', '18 ms',   '~10 row read; indexed by conversation_id.'],
       ['BM25 search (1,000 chunks)',            '2–4 ms',  '8 ms',    'Pure in-process JS — no I/O, no network call.'],
       ['System prompt assembly',               '< 1 ms',  '< 1 ms',  'String concatenation + JSON serialization.'],
       ['Gemma API call  ⚠ bottleneck',        '1.2 s',   '4.8 s',   'External HTTPS. Flash ~1.2 s, Pro ~3–5 s avg.'],
       ['SQLite message insert (2 rows)',        '1–2 ms',  '4 ms',    'WAL mode — synchronous write, immediate ACK.'],
       ['ElevenLabs TTS (200 words)',            '350 ms',  '800 ms',  'Turbo v2.5 model. Longer text → proportionally slower.'],
       ['Total local pipeline (excl. Gemma)',   '12–22 ms','38 ms',   'Steps 2–8 in the 12-step chat flow.'],
       ['End-to-end chat (Flash model)',        '1.2–2 s', '5.5 s',   'Dominated by Gemma API network round-trip.']],
      0.48, 1.52, 12.37, 5.35, GREEN, col_widths=[3.3, 1.25, 1.25, 6.57])
info_box(sl, 0.48, 6.65, 12.37, 0.62,
         '⚡  Optimisation Takeaway',
         '99% of end-to-end latency is the Gemma API call. Gemma Flash is 3–5× faster than Pro models. All local ops complete in under 25 ms total.',
         GREEN)
snum(sl, 59)

# ── SLIDE 60  REAL-WORLD DEPLOYMENTS ─────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Case Studies', 'Real-World Deployment Configurations', PURPLE)
configs60 = [
    ('🏥','Wellness Clinic Assistant',
     '· Model: Gemma 3 Flash\n· KB: Clinical guidelines PDFs\n· KB weight: 100% (strict)\n· Web search: disabled\n· Profile: full health fields\n· Fitness: Google Fit enabled\n· TTS: Rachel voice, Turbo v2.5', BLUE),
    ('🏢','HR Policy Bot',
     '· Model: Gemma 2.5 Flash\n· KB: Employee handbook, policies\n· KB weight: 90%\n· Own knowledge: 40%\n· Web search: disabled\n· Profile: name, department only\n· TTS: disabled', PURPLE),
    ('🎓','Student Tutor Bot',
     '· Model: Gemma Flash Preview\n· KB: Course PDFs, textbooks\n· KB weight: 70%\n· Own knowledge: 80%\n· Web search: 60%\n· Temperature: Creative (1.2)\n· TTS: Adam voice enabled', CYAN),
]
for i, (icon, title, desc, acc) in enumerate(configs60):
    R(sl, 0.48+i*4.3, 1.52, 4.05, 5.52, fill=_dim(acc,10), line=acc, lw=0.6, radius=True)
    R(sl, 0.48+i*4.3, 1.52, 0.06, 5.52, fill=acc)
    T(sl, icon+' '+title, 0.65+i*4.3, 1.62, 3.82, 0.36, size=12.5, bold=True, color=acc)
    hline(sl, 0.65+i*4.3, 2.06, 3.7, _dim(acc,4))
    T(sl, desc, 0.65+i*4.3, 2.20, 3.82, 4.55, size=11, color=MUTED)
snum(sl, 60)

# ── SLIDE 61  TRADE-OFFS ──────────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Trade-Offs', 'Honest Trade-Off Analysis', BLUE)
R(sl, 0.48, 1.52, 5.9, 5.58, fill=_dim(GREEN,11), line=_dim(GREEN,4), lw=0.5, radius=True)
R(sl, 0.48, 1.52, 0.06, 5.58, fill=GREEN)
T(sl, '✅  Where CIRA Excels', 0.68, 1.62, 5.7, 0.30, size=12.5, bold=True, color=GREEN)
bullets(sl, [
    ('~$12/month total fixed infra',     'SQLite + BM25 eliminates $70–250/month in vector DB and embedding API costs alone.'),
    ('Deep personalization built-in',    'Health profile + auto-extracted AI memories + live Google Fit data — all injected into every chat request automatically.'),
    ('Fast local pipeline: < 25 ms',    'JWT verify + DB reads + BM25 search + prompt assembly all complete in under 25 ms before the Gemma API call.'),
    ('Zero-code admin configuration',   'Model, temperature, system prompt, KB weights, TTS settings — all changeable from the admin UI with immediate effect.'),
    ('True multi-instance isolation',   'Separate OS processes, separate SQLite files, separate ports — Instance A crash cannot affect Instance B.'),
    ('Progressive enhancement',         'Core chat survives ElevenLabs outages, browser STT gaps, missing Google Fit — always shows text response.'),
    ('Fitness-aware AI (unique)',        'Injecting real-time wearable data (steps, HR, weight) into AI context is a unique capability in open-source chatbots.'),
    ('Model-agnostic architecture',     'Swap from Gemma Flash to Pro or Gemma 4 with one dropdown — same API contract, no code changes required.'),
], 0.68, 2.0, 5.6, 4.8, GREEN, size=9.5)
R(sl, 6.72, 1.52, 6.13, 5.58, fill=_dim(YELLOW,11), line=_dim(YELLOW,4), lw=0.5, radius=True)
R(sl, 6.72, 1.52, 0.06, 5.58, fill=YELLOW)
T(sl, '⚠️  Known Limitations', 6.92, 1.62, 5.9, 0.30, size=12.5, bold=True, color=YELLOW)
bullets(sl, [
    ('BM25 keyword gap',            'No synonym or semantic matching. "heart attack" won\'t retrieve a chunk containing "myocardial infarction".'),
    ('SQLite concurrency ceiling',  'WAL mode handles ~50 concurrent users per instance. > 1,000 concurrent users needs PostgreSQL migration.'),
    ('No response streaming',       'Full AI text returned at once — users see a loading indicator until the entire Gemma response is complete.'),
    ('Gemma API dependency',        'Single point of failure for AI responses. Gemma API outages make the chatbot unable to respond.'),
    ('Single-node SQLite',          'SQLite is file-based — cannot be shared across multiple hosts. Horizontal scaling requires DB migration.'),
    ('No rate limiting',            'No per-user request throttling. A runaway client can exhaust Gemma API quota or overload SQLite.'),
    ('Memory extraction cap',       'Memory extraction reads last 150 messages only — very old conversation context is not analysed.'),
    ('No refresh token rotation',   'JWT is 7-day non-rotatable. Stolen tokens remain valid until expiry with no revocation mechanism.'),
], 6.92, 2.0, 5.9, 4.8, YELLOW, size=9.5)
snum(sl, 61)

# ── SLIDE 62  FUTURE ROADMAP ──────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Roadmap', 'Future Roadmap', PURPLE)
T(sl, 'Near Term  (v2.0)', 0.48, 1.52, 5.9, 0.30, size=11.5, bold=True, color=BLUE)
for i, (title, desc, acc) in enumerate([
    ('Streaming Responses',   'Server-Sent Events (SSE) for real-time token streaming — eliminates perceived latency for long responses.',BLUE),
    ('Hybrid BM25 + Semantic','Add optional vector embeddings as opt-in alongside BM25 for queries requiring semantic understanding.',PURPLE),
    ('Rate Limiting',         'express-rate-limit middleware to prevent API abuse. Configurable per-user request limits in admin panel.',CYAN),
]):
    R(sl, 0.48, 1.88+i*1.44, 0.62, 0.62, fill=_dim(acc,4), line=acc, lw=0.6, radius=True)
    T(sl, str(i+1), 0.48, 1.88+i*1.44, 0.62, 0.62, size=15, bold=True, color=acc, align=PP_ALIGN.CENTER)
    R(sl, 1.2, 1.88+i*1.44, 5.18, 0.62, fill=_dim(acc,10), line=_dim(acc,4), lw=0.4, radius=True)
    T(sl, title, 1.34, 1.93+i*1.44, 4.9, 0.28, size=11.5, bold=True, color=WHITE)
    T(sl, desc,  1.34, 2.24+i*1.44, 4.9, 1.0,  size=9.5, color=MUTED)
T(sl, 'Long Term  (v3.0)', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=PURPLE)
for i, (title, desc, acc) in enumerate([
    ('PostgreSQL Migration',  'Abstract DB layer to support both SQLite and PostgreSQL for horizontal enterprise scaling.',PURPLE),
    ('Multi-Modal Support',   "Image uploads in chat using Gemma's vision capabilities — analyse charts, photos, documents inline.",GREEN),
    ('Wearable Integrations', 'Expand beyond Google Fit to Apple Health, Fitbit, Garmin APIs for broader fitness data coverage.',YELLOW),
]):
    R(sl, 6.72, 1.88+i*1.44, 0.62, 0.62, fill=_dim(acc,4), line=acc, lw=0.6, radius=True)
    T(sl, str(i+4), 6.72, 1.88+i*1.44, 0.62, 0.62, size=15, bold=True, color=acc, align=PP_ALIGN.CENTER)
    R(sl, 7.44, 1.88+i*1.44, 5.41, 0.62, fill=_dim(acc,10), line=_dim(acc,4), lw=0.4, radius=True)
    T(sl, title, 7.58, 1.93+i*1.44, 5.2, 0.28, size=11.5, bold=True, color=WHITE)
    T(sl, desc,  7.58, 2.24+i*1.44, 5.2, 1.0,  size=9.5, color=MUTED)
snum(sl, 62)

# ── SLIDE 63  DATA FLOW SUMMARY ───────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Data Flow', 'Complete Data Flow — One Chat Message', CYAN)
cols63 = [
    ('USER',       BLUE,   '1. Types message\n   ↓\n2. POST /api/chat\n   Bearer: JWT\n   ↓\n3. Waits for response\n   ↓\n4. Renders Markdown\n   ↓\n5. Plays TTS audio'),
    ('API SERVER', PURPLE, '1. Verifies JWT\n   ↓\n2. Reads user profile\n   ↓\n3. Reads memories\n   ↓\n4. Reads fitness data\n   ↓\n5. Reads 10 msgs\n   ↓\n6. Assembles prompt\n   ↓\n7. Saves both msgs'),
    ('RAG ENGINE', GREEN,  '1. Tokenize query\n   ↓\n2. Load all KB chunks\n   ↓\n3. BM25 score each\n   ↓\n4. Sort by relevance\n   ↓\n5. Return Top-5\n   ↓\n6. With doc names'),
    ('GEMINI API', YELLOW, '1. Receives context\n   ↓\n2. System instruction\n   ↓\n3. User profile data\n   ↓\n4. KB chunks\n   ↓\n5. Chat history\n   ↓\n6. Generates text\n   ↓\n7. Returns response'),
]
for i, (title, acc, content) in enumerate(cols63):
    R(sl, 0.42+i*3.24, 1.52, 3.06, 5.65, fill=_dim(acc,10), line=acc, lw=0.6, radius=True)
    R(sl, 0.42+i*3.24, 1.52, 3.06, 0.38, fill=_dim(acc,5))  # header fill
    T(sl, title, 0.50+i*3.24, 1.58, 2.9, 0.28, size=10.5, bold=True, color=acc, align=PP_ALIGN.CENTER)
    hline(sl, 0.52+i*3.24, 1.92, 2.84, _dim(acc,4))
    T(sl, content, 0.52+i*3.24, 2.04, 2.9, 4.88, size=10, color=MUTED, align=PP_ALIGN.CENTER)
snum(sl, 63)

# ── SLIDE 64  SECURITY CHECKLIST ─────────────────────────────────
sl = new_slide(prs, RGBColor(0x10, 0x06, 0x0b))
slide_header(sl, 'Security Checklist', 'Production Security Checklist', RED)
checks_pass = [
    'Passwords hashed with bcryptjs (10 salt rounds — ~100 ms hash time)',
    'JWT signed with HMAC-SHA256, 7-day expiry, role claim embedded',
    'CORS restricted to FRONTEND_URL environment variable only',
    'Admin routes double-protected: requireAuth + requireAdmin middleware chain',
    'OAuth2 read-only scopes for Google Fit (activity, body, heart_rate)',
    'OAuth2 tokens stored server-side in user.profile — never sent to frontend',
    'File upload type validation: MIME type + extension check (PDF/TXT/MD only)',
    'File size limit enforced at multer layer (20 MB maximum per document)',
    'Markdown rendered via react-markdown with HTML disabled — XSS prevented',
    'Error responses return {"error":"message"} — no stack traces exposed',
    'UUID-based IDs for all resources — no sequential ID enumeration attacks',
]
checks_warn = [
    'Rate limiting not yet implemented (express-rate-limit planned for v2.0)',
    'HTTPS enforcement delegated to Render.com / Netlify platform layer',
    'No maximum length validation on chat message input body',
    'No CSRF tokens (SPA + JWT Bearer header mitigates standard CSRF vectors)',
    'JWT secret rotation requires re-login for all active users (no refresh tokens)',
    'No audit logging for admin actions (settings changes, user deletions)',
]
T(sl, '✅  Implemented Security Controls', 0.48, 1.52, 5.9, 0.30, size=11.5, bold=True, color=GREEN)
for i, c in enumerate(checks_pass):
    check_row(sl, c, 0.48, 1.88+i*0.46, 5.9, passed=True)
T(sl, '⚠️  Known Gaps / Recommendations', 6.72, 1.52, 6.13, 0.30, size=11.5, bold=True, color=YELLOW)
for i, c in enumerate(checks_warn):
    check_row(sl, c, 6.72, 1.88+i*0.60, 6.13, passed=False)
snum(sl, 64)

# ── SLIDE 65  ADVANTAGES SUMMARY ─────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Advantages Summary', 'Why Choose CIRA — 8 Competitive Advantages', BLUE)
advs65 = [
    ('💰','Cost Efficient',      '~$12/month vs. $110–370 for equivalent stacks',BLUE),
    ('🧠','Deep Personalization','Profile + AI memories + live Google Fit data',PURPLE),
    ('🔍','Zero-Cost RAG',       'BM25: no vector DB or embedding API',CYAN),
    ('⚙️','Admin-First',         'Full bot config without any code changes',GREEN),
    ('🏗️','Multi-Instance',      'N isolated bots from one codebase',YELLOW),
    ('🤖','Model Agnostic',      '6 Gemma variants, switch without code',RED),
    ('🛡️','Secure by Design',    'JWT + bcrypt + RBAC + CORS + OAuth2',BLUE),
    ('📱','Rich UX',             'TTS · STT · Markdown · themes · responsive',PURPLE),
]
for i, (icon, title, desc, acc) in enumerate(advs65):
    c = i % 4; rr = i // 4
    R(sl, 0.42+c*3.24, 1.52+rr*2.62, 3.06, 2.48, fill=_dim(acc,9), line=acc, lw=0.6, radius=True)
    R(sl, 0.42+c*3.24, 1.52+rr*2.62, 0.06, 2.48, fill=acc)
    T(sl, icon, 0.58+c*3.24, 1.62+rr*2.62, 2.82, 0.50, size=20, align=PP_ALIGN.CENTER)
    T(sl, title, 0.58+c*3.24, 2.18+rr*2.62, 2.82, 0.30, size=11, bold=True, color=acc, align=PP_ALIGN.CENTER)
    T(sl, desc,  0.58+c*3.24, 2.52+rr*2.62, 2.82, 1.25, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)
snum(sl, 65)

# ── SLIDE 66  API DESIGN PRINCIPLES ──────────────────────────────
sl = new_slide(prs, BG2)
slide_header(sl, 'API Principles', 'REST API Design Principles', GREEN)
principles66 = [
    ('📐','Resource-Oriented URLs','URLs identify resources: /api/conversations/:id, /api/user/memories/:id. Actions expressed via HTTP verbs (GET/POST/PUT/PATCH/DELETE) — not verbs in URLs.',GREEN),
    ('🔄','Consistent Response Shape','Success: data payload JSON. Error: {"error":"message"} with correct HTTP code (400/401/403/500). Frontend handles both shapes uniformly.',BLUE),
    ('🔗','Stateless Requests','Every request carries full auth context (JWT Bearer token). No server session required. Enables load balancing without sticky sessions.',PURPLE),
    ('🏷️','Content Negotiation','JSON for all API responses. MP3 audio for /api/tts. Content-Type headers correctly set on all responses for proper client-side handling.',CYAN),
]
for i, (icon, title, desc, acc) in enumerate(principles66):
    c = i % 2; rr = i // 2
    R(sl, 0.48+c*6.48, 1.52+rr*2.72, 6.15, 2.56, fill=_dim(acc,10), line=acc, lw=0.5, radius=True)
    R(sl, 0.48+c*6.48, 1.52+rr*2.72, 0.06, 2.56, fill=acc)
    T(sl, icon+' '+title, 0.68+c*6.48, 1.64+rr*2.72, 5.75, 0.32, size=12.5, bold=True, color=acc)
    T(sl, desc, 0.68+c*6.48, 2.04+rr*2.72, 5.75, 1.85, size=10.5, color=MUTED)
snum(sl, 66)

# ── SLIDE 67  TECHNICAL INNOVATIONS ──────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Innovations', 'Technical Innovation Highlights', CYAN)
for i, (icon, title, desc, acc) in enumerate([
    ('🧬','AI-Powered Memory', 'Using Gemma to analyse conversation history and extract persistent user facts is a unique approach to implicit personalisation — no user survey or manual tagging required.', GREEN),
    ('⚖️','Weighted Knowledge Blending', 'Simultaneously weighting RAG, model knowledge, and web search is not a standard chatbot pattern. It enables fine-grained control over information sourcing per deployment.', PURPLE),
    ('🏃','Live Fitness Context', 'Injecting real-time wearable data (steps, HR, weight) into AI system prompts bridges the gap between fitness tracking and actionable AI advice — a first for open-source chatbots.', CYAN),
]):
    R(sl, 0.48+i*4.3, 1.52, 4.05, 5.52, fill=CARD, line=STROKE, lw=0.5, radius=True)
    R(sl, 0.48+i*4.3, 1.52, 0.06, 5.52, fill=acc)
    T(sl, icon, 0.65+i*4.3, 1.72, 3.82, 0.65, size=28, align=PP_ALIGN.CENTER)
    T(sl, title, 0.65+i*4.3, 2.48, 3.82, 0.40, size=13.5, bold=True, color=acc, align=PP_ALIGN.CENTER)
    hline(sl, 0.72+i*4.3, 2.96, 3.62, _dim(acc,4))
    T(sl, desc, 0.65+i*4.3, 3.14, 3.82, 2.7, size=11, color=MUTED, align=PP_ALIGN.CENTER)
snum(sl, 67)

# ── SLIDE 68  KEY TAKEAWAYS ───────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Key Takeaways', 'Key Takeaways', BLUE)
for i, (icon, title, desc, acc) in enumerate([
    ('🎯','Personalization is the core differentiator',
     'Health profile + auto-extracted AI memories + Google Fit live data injected into every Gemma request. No other open-source chatbot combines all three in one platform.',BLUE),
    ('💡','BM25 eliminates the biggest cost driver',
     'Zero embedding API calls + zero vector DB = $0 RAG infrastructure vs. $70–250/month for Pinecone/Weaviate. BM25 matches semantic search accuracy for domain-specific documents.',GREEN),
    ('⚡','Admin-first design enables non-technical management',
     'Gemma model, temperature, system prompt, RAG weights, TTS voice — all configurable via admin UI. Settings read from SQLite on every request; no server restart or redeploy ever needed.',PURPLE),
    ('🏗️','Architecture scales pragmatically',
     'SQLite WAL mode handles ~50 concurrent users per instance. Multiple instances scale horizontally. Clearly defined migration path to PostgreSQL for enterprise requirements.',YELLOW),
    ('🌱','Production-ready in under 30 minutes',
     'Clone repo → npm install → set GEMMA_API_KEY + JWT_SECRET → node start-instances.js. Deploy to Render + Netlify with render.yaml and netlify.toml already configured.',CYAN),
]):
    R(sl, 0.48, 1.52+i*1.12, 12.37, 1.0, fill=_dim(acc,10), line=acc, lw=0.5, radius=True)
    R(sl, 0.48, 1.52+i*1.12, 0.06, 1.0, fill=acc)
    T(sl, icon, 0.68, 1.60+i*1.12, 0.65, 0.65, size=18, align=PP_ALIGN.CENTER)
    T(sl, title, 1.45, 1.62+i*1.12, 6.8, 0.34, size=12.5, bold=True, color=WHITE)
    T(sl, desc,  1.45, 2.0+i*1.12,  10.8, 0.50, size=10.5, color=MUTED)
snum(sl, 68)

# ── SLIDE 69  DEMO HIGHLIGHTS ─────────────────────────────────────
sl = new_slide(prs)
slide_header(sl, 'Live Demo', 'What to Look For in the Demo', GREEN)
demo69 = [
    ('💬','Multi-Conversation',  'Create two separate conversations. Notice each has an auto-generated title and independent message history.',GREEN),
    ('📚','RAG in Action',       'Ask about content from an uploaded document. Notice source citation references appear alongside the AI response.',CYAN),
    ('🔊','TTS + STT',           'Click the TTS button to hear the response in ElevenLabs voice. Use the microphone to speak a question.',BLUE),
    ('👤','Profile Impact',      'Fill the health profile. Ask a fitness question — notice how the answer references your specific profile data.',PURPLE),
    ('💾','Memory Extraction',   'After several conversations, trigger memory extraction. Watch Gemma auto-identify learned facts about you.',YELLOW),
    ('⚙️','Admin Settings',      'Switch the AI model, change the system prompt, upload a document — all without restarting the server.',RED),
]
for i, (icon, title, desc, acc) in enumerate(demo69):
    c = i % 3; rr = i // 3
    card(sl, 0.48+c*4.3, 1.52+rr*2.72, 4.05, 2.55, icon, title, desc, acc)
snum(sl, 69)

# ── SLIDE 70  THANK YOU / Q&A ─────────────────────────────────────
sl = new_slide(prs, RGBColor(0x06, 0x08, 0x18))
R(sl, 0, 0, 13.333, 1.35, fill=RGBColor(0x08, 0x0e, 0x2a))   # top bar
R(sl, 0, 6.15, 13.333, 1.35, fill=RGBColor(0x08, 0x0e, 0x2a)) # bottom bar
R(sl, 0, 0, 0.18, 7.5, fill=BLUE)                              # left accent
R(sl, 13.153, 0, 0.18, 7.5, fill=PURPLE)                       # right accent
badge(sl, 'Thank You', 5.82, 0.28, BLUE)
T(sl, 'Questions?', 0.8, 1.52, 11.73, 1.65, size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(sl, 'CIRA — Intelligent Conversational AI',
  0.8, 3.18, 11.73, 0.72, size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
hline(sl, 2.8, 4.0, 7.73, DIM)
for i, (icon, lbl, val, acc) in enumerate([
    ('🔗','GitHub Repository', 'bashab18/chatbot', BLUE),
    ('⚡','Powered By',        'Google Gemma API', PURPLE),
    ('🚀','Deploy On',         'Render + Netlify',  CYAN),
]):
    R(sl, 2.62+i*2.95, 4.22, 2.7, 1.52, fill=_dim(acc,9), line=acc, lw=0.6, radius=True)
    T(sl, icon, 2.62+i*2.95, 4.30, 2.7, 0.48, size=18, align=PP_ALIGN.CENTER)
    T(sl, lbl,  2.62+i*2.95, 4.78, 2.7, 0.28, size=9,  color=MUTED, align=PP_ALIGN.CENTER)
    T(sl, val,  2.62+i*2.95, 5.06, 2.7, 0.42, size=12, bold=True, color=acc, align=PP_ALIGN.CENTER)
tags70 = [('React 18',BLUE),('Node.js',PURPLE),('SQLite WAL',GREEN),('BM25 RAG',CYAN),
          ('JWT Auth',YELLOW),('Google Fit',RED),('ElevenLabs TTS',BLUE),('Multi-Instance',PURPLE)]
x70 = 0.48
for tag, col in tags70:
    badge(sl, tag, x70, 6.52, col); x70 += len(tag)*0.107 + 0.72
snum(sl, 70)

# ════════════════════════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), 'CIRA_Chatbot_Presentation.pptx')
prs.save(out)
print(f'✅  Saved  →  {out}')
print(f'   Slides : {len(prs.slides)}')
